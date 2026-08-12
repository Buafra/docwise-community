import base64
import hashlib
import json
import os
import re
import shutil
import sqlite3
import subprocess
import sys
import unicodedata
import tempfile
import threading
import time
from datetime import datetime
from pathlib import Path
from typing import Optional

from fastapi import FastAPI, File, Form, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse, Response
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel

try:
    import fitz  # PyMuPDF
    FITZ_IMPORT_ERROR = None
except Exception as _fitz_exc:
    fitz = None
    FITZ_IMPORT_ERROR = str(_fitz_exc)

try:
    from PIL import Image
except Exception:
    Image = None

try:
    import pytesseract
except Exception:
    pytesseract = None

try:
    from openai import OpenAI
except Exception:
    OpenAI = None

try:
    from docx import Document
except Exception:
    Document = None

try:
    from openpyxl import Workbook, load_workbook
except Exception:
    load_workbook = None
    Workbook = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

ROOT = Path(__file__).resolve().parent
DATA = ROOT / "data"
UPLOADS = DATA / "uploads"
ARCHIVE = DATA / "archive"
DB_PATH = DATA / "docwise.sqlite3"
STATIC = ROOT / "static_astryx" if (ROOT / "static_astryx").exists() else ROOT / "static"


def load_env_file() -> None:
    """Load KEY=VALUE lines from .env next to app.py, without overriding
    variables already set in the environment. Keeps API keys out of the
    launchers and out of git (.env is ignored)."""
    env_path = ROOT / ".env"
    if not env_path.exists():
        return
    try:
        for line in env_path.read_text(encoding="utf-8").splitlines():
            line = line.strip()
            if not line or line.startswith("#") or "=" not in line:
                continue
            key, _, value = line.partition("=")
            key = key.strip()
            value = value.strip().strip('"').strip("'")
            if key and key not in os.environ:
                os.environ[key] = value
    except Exception:
        pass


load_env_file()

for folder in (DATA, UPLOADS, ARCHIVE):
    folder.mkdir(parents=True, exist_ok=True)

SUPPORTED = {
    ".pdf", ".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff",
    ".txt", ".md", ".csv", ".json", ".rtf",
    ".docx", ".docm", ".dotx", ".doc",
    ".xlsx", ".xlsm", ".xltx", ".xls",
    ".pptx", ".pptm", ".ppsx", ".ppt",
}
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff"}
TEXT_EXTS = {".txt", ".md", ".csv", ".json"}
WORD_EXTS = {".docx", ".docm", ".dotx"}
EXCEL_EXTS = {".xlsx", ".xlsm", ".xltx"}
POWERPOINT_EXTS = {".pptx", ".pptm", ".ppsx"}
LEGACY_OFFICE_EXTS = {".doc", ".xls", ".ppt"}
ARABIC_RE = re.compile(r"[\u0600-\u06FF]")
TASHKEEL_RE = re.compile(r"[\u0610-\u061A\u064B-\u065F\u0670\u06D6-\u06ED]")
ARABIC_INDIC_DIGITS = str.maketrans("\u0660\u0661\u0662\u0663\u0664\u0665\u0666\u0667\u0668\u0669\u06F0\u06F1\u06F2\u06F3\u06F4\u06F5\u06F6\u06F7\u06F8\u06F9\u066B\u066C", "01234567890123456789.,")
CURRENCY_TOKEN = r"(?:AED|SAR|USD|EUR|GBP|QAR|KWD|BHD|OMR|EGP|JOD|\$|\u20AC|\u00A3|\u062F\.\u0625|\u0631\.\u0633|\u0631\.\u0642|\u062F\.\u0643|\u062F\.\u0628|\u0631\.\u0639|\u062F\.\u0627|\u062F\u0631\u0647\u0645|\u0631\u064A\u0627\u0644|\u062F\u064A\u0646\u0627\u0631|\u062C\u0646\u064A\u0647|\u062F\u0648\u0644\u0627\u0631|\u064A\u0648\u0631\u0648)"
AMOUNT_NUM = r"[0-9][0-9,]*(?:\.\d{1,2})?"


def normalize_digits(text: str) -> str:
    """Arabic-Indic digits to ASCII so amount/date extraction sees \u0669\u0665\u0660 as 950,
    with NFKC folding presentation-form glyphs into real Arabic letters."""
    return unicodedata.normalize("NFKC", text or "").translate(ARABIC_INDIC_DIGITS)

app = FastAPI(title="DocWise Archive AI", version="0.1.0")
app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_methods=["*"], allow_headers=["*"])
DB_LOCK = threading.Lock()
SCAN_STATE = {"running": False, "last": None, "message": "Ready", "indexed": 0, "errors": 0}


def now_iso() -> str:
    return datetime.now().isoformat(timespec="seconds")


def db() -> sqlite3.Connection:
    conn = sqlite3.connect(DB_PATH, check_same_thread=False)
    conn.row_factory = sqlite3.Row
    return conn


def init_db() -> None:
    with DB_LOCK, db() as conn:
        conn.executescript(
            """
            CREATE TABLE IF NOT EXISTS documents (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                path TEXT NOT NULL UNIQUE,
                original_name TEXT,
                source_type TEXT NOT NULL DEFAULT 'folder',
                file_ext TEXT,
                size INTEGER DEFAULT 0,
                mtime REAL DEFAULT 0,
                sha256 TEXT,
                title TEXT,
                doc_type TEXT DEFAULT 'general',
                language TEXT DEFAULT 'unknown',
                date_guess TEXT,
                company TEXT,
                amount TEXT,
                tags TEXT DEFAULT '[]',
                summary TEXT,
                fields TEXT DEFAULT '{}',
                text TEXT DEFAULT '',
                normalized_text TEXT DEFAULT '',
                status TEXT DEFAULT 'new',
                ocr_engine TEXT DEFAULT 'none',
                ocr_quality TEXT DEFAULT 'unknown',
                ocr_score REAL DEFAULT 0,
                error TEXT,
                created_at TEXT,
                updated_at TEXT
            );
            CREATE TABLE IF NOT EXISTS chunks (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                document_id INTEGER NOT NULL,
                page INTEGER DEFAULT 1,
                chunk_index INTEGER DEFAULT 0,
                text TEXT NOT NULL,
                normalized_text TEXT NOT NULL,
                token_count INTEGER DEFAULT 0,
                FOREIGN KEY(document_id) REFERENCES documents(id) ON DELETE CASCADE
            );
            CREATE TABLE IF NOT EXISTS chunk_embeddings (
                chunk_id INTEGER PRIMARY KEY,
                document_id INTEGER NOT NULL,
                model TEXT NOT NULL,
                dims INTEGER NOT NULL,
                embedding TEXT NOT NULL,
                created_at TEXT,
                FOREIGN KEY(chunk_id) REFERENCES chunks(id) ON DELETE CASCADE
            );
            CREATE VIRTUAL TABLE IF NOT EXISTS chunk_fts USING fts5(
                chunk_id UNINDEXED,
                document_id UNINDEXED,
                title,
                text,
                normalized_text,
                tokenize='unicode61'
            );
            CREATE TABLE IF NOT EXISTS page_words (
                document_id INTEGER NOT NULL,
                page INTEGER NOT NULL,
                words TEXT NOT NULL,
                PRIMARY KEY(document_id, page),
                FOREIGN KEY(document_id) REFERENCES documents(id) ON DELETE CASCADE
            );
            CREATE TABLE IF NOT EXISTS folders (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                path TEXT NOT NULL UNIQUE,
                recursive INTEGER DEFAULT 1,
                watch INTEGER DEFAULT 1,
                last_scan TEXT,
                created_at TEXT
            );
            CREATE TABLE IF NOT EXISTS eval_cases (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                question TEXT NOT NULL,
                expected TEXT DEFAULT '',
                must_cite TEXT DEFAULT '',
                created_at TEXT
            );
            CREATE INDEX IF NOT EXISTS idx_documents_type ON documents(doc_type);
            CREATE INDEX IF NOT EXISTS idx_documents_lang ON documents(language);
            CREATE INDEX IF NOT EXISTS idx_chunks_document ON chunks(document_id);
            CREATE INDEX IF NOT EXISTS idx_embeddings_document ON chunk_embeddings(document_id);
            """
        )
        for stmt in (
            "ALTER TABLE chunks ADD COLUMN token_count INTEGER DEFAULT 0",
            "ALTER TABLE documents ADD COLUMN fields TEXT DEFAULT '{}'",
            "ALTER TABLE documents ADD COLUMN ocr_quality TEXT DEFAULT 'unknown'",
            "ALTER TABLE documents ADD COLUMN ocr_score REAL DEFAULT 0",
        ):
            try:
                conn.execute(stmt)
            except sqlite3.OperationalError:
                pass



init_db()


def normalize_arabic(text: str) -> str:
    # NFKC folds Arabic presentation-form ligature glyphs (U+FB50-U+FEFF) back
    # to real letters - broken embedded text layers in Arabic PDFs are full of
    # them, and keyword matching is blind to them otherwise.
    text = unicodedata.normalize("NFKC", text or "")
    text = TASHKEEL_RE.sub("", text)
    text = re.sub("[إأآٱ]", "ا", text)
    text = text.replace("ى", "ي").replace("ؤ", "و").replace("ئ", "ي")
    text = text.replace("ة", "ه")
    text = re.sub(r"[؟،؛ـ]", " ", text)
    text = text.lower()
    text = re.sub(r"[^\w\u0600-\u06FF]+", " ", text, flags=re.UNICODE)
    return re.sub(r"\s+", " ", text).strip()


def safe_json_load(value: str, default):
    try:
        return json.loads(value or "")
    except Exception:
        return default


def contains_keyword(normalized_text: str, keyword: str) -> bool:
    kw = normalize_arabic(keyword)
    if not kw:
        return False
    padded = f" {normalized_text} "
    if " " in kw:
        return f" {kw} " in padded
    if re.fullmatch(r"[a-z0-9]+", kw):
        return f" {kw} " in padded
    # Arabic classification keywords must match word-like tokens, not substrings.
    # This prevents false matches like عقد inside المعقد, or طبي inside unrelated words.
    words = set(normalized_text.split())
    variants = {kw, f"ال{kw}", f"و{kw}", f"بال{kw}", f"لل{kw}"}
    for word in words:
        if word in variants:
            return True
        stripped = word
        for prefix in ("وال", "بال", "لل", "ال", "و", "ب", "ل"):
            if stripped.startswith(prefix) and len(stripped) > len(prefix) + 1:
                stripped = stripped[len(prefix):]
                break
        if stripped == kw:
            return True
    return False


def approx_token_count(text: str) -> int:
    # Good enough for chunk sizing/index stats without adding tokenizer dependencies.
    words = re.findall(r"[\w\u0600-\u06FF]+", text or "", flags=re.UNICODE)
    return max(1, int(len(words) * 1.25)) if words else 0


def cosine_similarity(a: list[float], b: list[float]) -> float:
    if not a or not b or len(a) != len(b):
        return 0.0
    dot = sum(x * y for x, y in zip(a, b))
    na = sum(x * x for x in a) ** 0.5
    nb = sum(y * y for y in b) ** 0.5
    return dot / (na * nb) if na and nb else 0.0


def local_embedding(text: str, dims: int = 384) -> list[float]:
    """Deterministic multilingual lexical vector fallback when no embedding API is configured."""
    vec = [0.0] * dims
    norm = normalize_arabic(text)
    tokens = [t for t in norm.split() if len(t) > 1]
    features = []
    for tok in tokens:
        features.append(tok)
        if len(tok) >= 4:
            features.extend(tok[i:i+3] for i in range(len(tok) - 2))
    for feat in features:
        h = int(hashlib.sha256(feat.encode("utf-8")).hexdigest()[:12], 16)
        idx = h % dims
        sign = 1.0 if (h >> 3) & 1 else -1.0
        vec[idx] += sign
    length = sum(v * v for v in vec) ** 0.5
    return [v / length for v in vec] if length else vec


# --- Local semantic embeddings: multilingual-e5-small via ONNX (no PyTorch) ---
try:
    import numpy as _np
    import onnxruntime as _ort
    from tokenizers import Tokenizer as _Tokenizer
except Exception:
    _np = _ort = _Tokenizer = None

MODELS_DIR = ROOT / "models" / "multilingual-e5-small"
E5_MODEL_FILE = MODELS_DIR / "model_quantized.onnx"
E5_TOKENIZER_FILE = MODELS_DIR / "tokenizer.json"
E5_DOWNLOADS = [
    (E5_MODEL_FILE, "https://huggingface.co/Xenova/multilingual-e5-small/resolve/main/onnx/model_quantized.onnx"),
    (E5_TOKENIZER_FILE, "https://huggingface.co/Xenova/multilingual-e5-small/resolve/main/tokenizer.json"),
]
ONNX_EMBED_MODEL = "local-e5-small-q8-v1"
HASH_EMBED_MODEL = "local-hash-multilingual-v1"
EMBED_STATE = {
    "runtime": bool(_ort and _Tokenizer and _np),
    "ready": False,
    "downloading": False,
    "progress": "",
    "backfill_total": 0,
    "backfill_done": 0,
}
_E5_LOCK = threading.Lock()
_E5: dict = {"session": None, "tokenizer": None, "needs_token_type": False}


def _e5_files_ready() -> bool:
    try:
        return (
            E5_MODEL_FILE.exists() and E5_MODEL_FILE.stat().st_size > 50_000_000
            and E5_TOKENIZER_FILE.exists() and E5_TOKENIZER_FILE.stat().st_size > 1_000_000
        )
    except OSError:
        return False


def _download_e5_models() -> None:
    import urllib.request

    EMBED_STATE["downloading"] = True
    try:
        MODELS_DIR.mkdir(parents=True, exist_ok=True)
        for dest, url in E5_DOWNLOADS:
            if dest.exists() and dest.stat().st_size > 1_000_000:
                continue
            tmp = dest.with_suffix(dest.suffix + ".part")
            with urllib.request.urlopen(url, timeout=60) as resp, tmp.open("wb") as out:
                total = int(resp.headers.get("Content-Length") or 0)
                got = 0
                while True:
                    block = resp.read(1024 * 512)
                    if not block:
                        break
                    out.write(block)
                    got += len(block)
                    mb = got // (1024 * 1024)
                    EMBED_STATE["progress"] = f"downloading {dest.name}: {mb}MB" + (f"/{total // (1024 * 1024)}MB" if total else "")
            tmp.replace(dest)
        EMBED_STATE["progress"] = "semantic model downloaded"
    except Exception as exc:
        EMBED_STATE["progress"] = f"model download failed (search still works, retrying later): {str(exc)[:120]}"
    finally:
        EMBED_STATE["downloading"] = False


def _load_e5() -> Optional[dict]:
    if not EMBED_STATE["runtime"]:
        return None
    with _E5_LOCK:
        if _E5["session"] is not None:
            return _E5
        if not _e5_files_ready():
            return None
        try:
            tok = _Tokenizer.from_file(str(E5_TOKENIZER_FILE))
            tok.enable_truncation(max_length=512)
            sess = _ort.InferenceSession(str(E5_MODEL_FILE), providers=["CPUExecutionProvider"])
            _E5.update({
                "session": sess,
                "tokenizer": tok,
                "needs_token_type": any(i.name == "token_type_ids" for i in sess.get_inputs()),
            })
            EMBED_STATE.update({"ready": True, "progress": "semantic embeddings active"})
            return _E5
        except Exception as exc:
            EMBED_STATE["progress"] = f"model load failed: {str(exc)[:120]}"
            return None


def onnx_embed_batch(texts: list[str], kind: str = "passage") -> Optional[list[list[float]]]:
    """e5 models expect 'query: ' / 'passage: ' prefixes; mean-pool + L2 normalize."""
    e5 = _load_e5()
    if not e5:
        return None
    prefixed = [f"{kind}: {(t or ' ')[:4000]}" for t in texts]
    encs = [e5["tokenizer"].encode(t) for t in prefixed]
    maxlen = max(1, max(len(e.ids) for e in encs))
    ids = _np.array([e.ids + [1] * (maxlen - len(e.ids)) for e in encs], dtype=_np.int64)
    mask = _np.array([e.attention_mask + [0] * (maxlen - len(e.attention_mask)) for e in encs], dtype=_np.int64)
    feed = {"input_ids": ids, "attention_mask": mask}
    if e5["needs_token_type"]:
        feed["token_type_ids"] = _np.zeros_like(ids)
    out = e5["session"].run(None, feed)[0]
    mask_f = mask[:, :, None].astype(_np.float32)
    emb = (out * mask_f).sum(axis=1) / _np.clip(mask_f.sum(axis=1), 1e-9, None)
    emb = emb / _np.clip(_np.linalg.norm(emb, axis=1, keepdims=True), 1e-9, None)
    return [row.astype(float).tolist() for row in emb]


def embed_text(text: str, kind: str = "passage") -> tuple[list[float], str]:
    mode = os.environ.get("DOCWISE_EMBEDDINGS", "auto").lower()
    if mode not in ("local", "onnx") and os.environ.get("OPENAI_API_KEY") and OpenAI:
        try:
            client = OpenAI()
            model = os.environ.get("DOCWISE_EMBED_MODEL", "text-embedding-3-small")
            resp = client.embeddings.create(model=model, input=(text or "")[:8000])
            return list(resp.data[0].embedding), model
        except Exception:
            if mode == "openai":
                raise
    if mode != "hash":
        vecs = onnx_embed_batch([text], kind=kind)
        if vecs:
            return vecs[0], ONNX_EMBED_MODEL
    return local_embedding(text), HASH_EMBED_MODEL


def active_embedding_model() -> str:
    mode = os.environ.get("DOCWISE_EMBEDDINGS", "auto").lower()
    if mode not in ("local", "onnx") and os.environ.get("OPENAI_API_KEY") and OpenAI:
        return os.environ.get("DOCWISE_EMBED_MODEL", "text-embedding-3-small")
    if mode != "hash" and EMBED_STATE["ready"]:
        return ONNX_EMBED_MODEL
    return HASH_EMBED_MODEL


def embedding_maintenance_loop() -> None:
    """Download the semantic model when missing, then upgrade old hash vectors."""
    time.sleep(5)
    while True:
        try:
            if EMBED_STATE["runtime"] and not _e5_files_ready():
                _download_e5_models()
                time.sleep(10)
                continue
            _load_e5()
            if active_embedding_model() != ONNX_EMBED_MODEL:
                time.sleep(60)
                continue
            with DB_LOCK, db() as conn:
                total = conn.execute(
                    "SELECT COUNT(*) c FROM chunk_embeddings WHERE model=?", (HASH_EMBED_MODEL,)
                ).fetchone()["c"]
                rows = conn.execute(
                    """
                    SELECT ce.chunk_id, ce.document_id, ch.text, d.title
                    FROM chunk_embeddings ce
                    JOIN chunks ch ON ch.id=ce.chunk_id
                    JOIN documents d ON d.id=ce.document_id
                    WHERE ce.model=? LIMIT 32
                    """,
                    (HASH_EMBED_MODEL,),
                ).fetchall()
            if not rows:
                EMBED_STATE.update({"backfill_total": 0, "backfill_done": 0})
                time.sleep(30)
                continue
            EMBED_STATE["backfill_total"] = max(EMBED_STATE["backfill_total"], total)
            vecs = onnx_embed_batch([f"{r['title']}\n{r['text']}" for r in rows], kind="passage")
            if not vecs:
                time.sleep(60)
                continue
            with DB_LOCK, db() as conn:
                for r, v in zip(rows, vecs):
                    conn.execute(
                        "UPDATE chunk_embeddings SET model=?, dims=?, embedding=?, created_at=? WHERE chunk_id=?",
                        (ONNX_EMBED_MODEL, len(v), json.dumps(v), now_iso(), r["chunk_id"]),
                    )
            EMBED_STATE["backfill_done"] += len(rows)
            time.sleep(0.5)
        except Exception:
            time.sleep(60)


def fts_query_from_terms(terms: list[str]) -> str:
    clean = []
    for term in terms[:12]:
        term = normalize_arabic(term)
        if len(term) < 2:
            continue
        term = re.sub(r'[^\w\u0600-\u06FF]+', ' ', term, flags=re.UNICODE).strip()
        if term:
            clean.append('"' + term.replace('"', '') + '"')
    return " OR ".join(dict.fromkeys(clean))


def query_plan(question: str) -> dict:
    q = normalize_arabic(question)
    plan = {"doc_type": None, "field": None, "sort": None, "terms": query_terms(question)}
    doc_rules = {
        "invoice": ["فاتوره", "الفاتوره", "bill", "invoice", "كهرباء", "utility"],
        "contract": ["عقد", "ايجار", "contract", "lease", "rent"],
        "id": ["هويه", "جواز", "passport", "emirates id"],
        "receipt": ["ايصال", "receipt", "payment"],
        "bank": ["بنك", "bank", "iban", "statement"],
        "medical": ["طبي", "مستشفي", "medical", "hospital", "clinic"],
    }
    for doc_type, kws in doc_rules.items():
        if any(contains_keyword(q, kw) for kw in kws):
            plan["doc_type"] = doc_type
            break
    if any(contains_keyword(q, kw) for kw in ["كم", "قيمه", "مبلغ", "total", "amount", "price"]):
        plan["field"] = "amount"
    if any(contains_keyword(q, kw) for kw in ["تاريخ", "متي", "date", "expiry", "انتهاء", "due"]):
        plan["field"] = plan["field"] or "date"
    if any(contains_keyword(q, kw) for kw in ["اخر", "احدث", "latest", "newest", "recent"]):
        plan["sort"] = "latest"
    return plan


def delete_chunk_indexes_for_doc(conn: sqlite3.Connection, doc_id: int) -> None:
    chunk_ids = [r["id"] for r in conn.execute("SELECT id FROM chunks WHERE document_id=?", (doc_id,)).fetchall()]
    if chunk_ids:
        placeholders = ",".join("?" for _ in chunk_ids)
        conn.execute(f"DELETE FROM chunk_embeddings WHERE chunk_id IN ({placeholders})", chunk_ids)
    conn.execute("DELETE FROM chunk_fts WHERE document_id=?", (doc_id,))


def add_chunk_indexes(conn: sqlite3.Connection, chunk_id: int, doc_id: int, title: str, text: str, normalized_text: str) -> None:
    conn.execute(
        "INSERT INTO chunk_fts(chunk_id, document_id, title, text, normalized_text) VALUES(?,?,?,?,?)",
        (chunk_id, doc_id, title or "", text or "", normalized_text or ""),
    )
    try:
        emb, model = embed_text(f"{title}\n{text}")
        conn.execute(
            "INSERT OR REPLACE INTO chunk_embeddings(chunk_id, document_id, model, dims, embedding, created_at) VALUES(?,?,?,?,?,?)",
            (chunk_id, doc_id, model, len(emb), json.dumps(emb), now_iso()),
        )
    except Exception as exc:
        # FTS still works if embeddings fail.
        conn.execute(
            "INSERT OR REPLACE INTO chunk_embeddings(chunk_id, document_id, model, dims, embedding, created_at) VALUES(?,?,?,?,?,?)",
            (chunk_id, doc_id, f"embedding-error:{str(exc)[:80]}", 0, "[]", now_iso()),
        )


def file_hash(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as f:
        for chunk in iter(lambda: f.read(1024 * 1024), b""):
            h.update(chunk)
    return h.hexdigest()


def has_arabic(text: str) -> bool:
    return bool(ARABIC_RE.search(text or ""))


def text_is_garbled(text: str) -> bool:
    """Detect PDF text extraction that is technically present but useless, e.g. Arabic rendered as ?????."""
    clean = re.sub(r"\s+", "", text or "")
    if not clean:
        return True
    bad_chars = clean.count("?") + clean.count("�") + clean.count("□") + clean.count("■")
    weird_chars = sum(1 for ch in clean if not re.match(r"[A-Za-z0-9\u0600-\u06FF\s.,:;!؟،؛()\[\]{}_/\\+\-=٪%$€£¥@#&*'\"|<>\n\r\t-]", ch))
    greek_garbage = len(re.findall(r"[\u0370-\u03FF\u1F00-\u1FFF]", clean))
    if len(clean) >= 8 and bad_chars / max(1, len(clean)) > 0.22:
        return True
    if greek_garbage >= 2:
        return True
    if len(clean) >= 20 and weird_chars / max(1, len(clean)) > 0.08:
        return True
    if re.fullmatch(r"[?\W_\d]+", clean) and clean.count("?") >= 4:
        return True
    return False


def should_ocr_pdf_page(text: str) -> bool:
    stripped = (text or "").strip()
    return len(stripped) < 25 or text_is_garbled(stripped)


def tesseract_cmd() -> Optional[str]:
    env = os.environ.get("TESSERACT_CMD")
    if env and Path(env).exists():
        return env
    found = shutil.which("tesseract")
    if found:
        return found
    candidates = [
        r"C:\Program Files\Tesseract-OCR\tesseract.exe",
        r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
    ]
    for c in candidates:
        if Path(c).exists():
            return c
    return None


def tessdata_dir() -> Optional[str]:
    env = os.environ.get("TESSDATA_PREFIX")
    if env and Path(env).exists():
        return env
    local = ROOT / "tessdata"
    if local.exists():
        return str(local)
    cmd = tesseract_cmd()
    if cmd:
        installed = Path(cmd).parent / "tessdata"
        if installed.exists():
            return str(installed)
    return None


def available_languages() -> list[str]:
    folder = tessdata_dir()
    if not folder:
        return []
    return sorted([p.stem for p in Path(folder).glob("*.traineddata")])


def available_ocr() -> dict:
    cmd = tesseract_cmd()
    langs = available_languages()
    return {
        "tesseract": bool(cmd and pytesseract and Image),
        "tesseract_cmd": cmd,
        "tessdata_dir": tessdata_dir(),
        "languages": langs,
        "arabic": "ara" in langs,
        "english": "eng" in langs,
        "office": {"word": bool(Document), "excel": bool(load_workbook), "powerpoint": bool(Presentation)},
        "engine_mode": os.environ.get("DOCWISE_OCR", "auto").lower(),
        "embedding_mode": os.environ.get("DOCWISE_EMBEDDINGS", "auto"),
        "embedding_model": active_embedding_model(),
        "embedding": dict(EMBED_STATE),
        "ollama": ollama_status(),
        "azure": bool(os.environ.get("AZURE_DI_ENDPOINT") and os.environ.get("AZURE_DI_KEY")),
        "openai_vision": bool(os.environ.get("OPENAI_API_KEY") and OpenAI),
        "openai_text": bool(os.environ.get("OPENAI_API_KEY") and OpenAI),
    }


def ocr_text_score(text: str) -> int:
    clean = (text or "").strip()
    if not clean:
        return 0
    arabic = len(re.findall(r"[\u0600-\u06FF]", clean))
    latin = len(re.findall(r"[A-Za-z]", clean))
    digits = len(re.findall(r"\d", clean))
    words = len(re.findall(r"[\w\u0600-\u06FF]{2,}", clean, flags=re.UNICODE))
    bad = clean.count("?") + clean.count("�") + clean.count("□")
    weird = sum(1 for ch in clean if not re.match(r"[A-Za-z0-9\u0600-\u06FF\s.,:;!؟،؛()\[\]{}_/\\+\-=٪%$€£¥@#&*'\"|<>\n\r\t-]", ch))
    # Arabic and Latin weighted equally: overweighting Arabic makes the
    # Arabic-only pass win on English documents with garbage glyph output.
    return arabic * 2 + latin * 2 + digits + words * 3 - bad * 10 - weird * 2


def ocr_quality(text: str) -> tuple[str, float]:
    clean = (text or "").strip()
    if not clean:
        return "empty", 0.0
    score = ocr_text_score(clean)
    chars = max(1, len(clean))
    bad = clean.count("?") + clean.count("�") + clean.count("□")
    bad_ratio = bad / chars
    word_count = len(re.findall(r"[\w\u0600-\u06FF]{2,}", clean, flags=re.UNICODE))
    density = score / chars
    if bad_ratio > 0.08 or word_count < 3:
        return "poor", round(max(0.0, density), 3)
    if density < 0.55:
        return "poor", round(density, 3)
    if density < 1.15:
        return "medium", round(density, 3)
    return "good", round(density, 3)


def should_try_vision_fallback(text: str, engine: str = "") -> bool:
    quality, _ = ocr_quality(text)
    if quality in ("empty", "poor"):
        return True
    # Tesseract's own confidence catches rubbish that merely *counts* like good
    # text: ID cards, passports, stamps, busy security backgrounds.
    m = re.search(r":conf(\d+)", engine or "")
    if m and int(m.group(1)) < 65:
        return True
    # Arabic Tesseract with diacritics often looks plausible but wrong; allow opt-in automatic cloud fallback.
    return bool(os.environ.get("DOCWISE_VISION_FALLBACK", "1") == "1" and "tesseract" in engine and has_arabic(text) and quality == "medium")


def ocr_image_variants(img):
    try:
        from PIL import ImageOps, ImageFilter
        variants = [("orig", img)]
        gray = ImageOps.grayscale(img)
        variants.append(("gray", gray))
        scale = 2 if max(img.size) < 2200 else 1
        if scale > 1:
            big = gray.resize((gray.width * scale, gray.height * scale))
            variants.append(("gray2x", big))
            variants.append(("sharp2x", big.filter(ImageFilter.SHARPEN)))
            variants.append(("threshold2x", big.point(lambda x: 255 if x > 180 else 0)))
        return variants
    except Exception:
        return [("orig", img)]


def ocr_image_tesseract(path: Path) -> tuple[str, str, Optional[dict]]:
    cmd = tesseract_cmd()
    if not (cmd and pytesseract and Image):
        return "", "tesseract-not-installed"
    pytesseract.pytesseract.tesseract_cmd = cmd
    img = Image.open(path)
    errors = []
    langs = set(available_languages())
    attempts = []
    if {"ara", "eng"}.issubset(langs):
        attempts.extend(["ara+eng", "eng+ara"])
    if "ara" in langs:
        attempts.append("ara")
    if "eng" in langs:
        attempts.append("eng")
    attempts = attempts or ["eng"]
    # Pass tessdata location via env, not --tessdata-dir: an unquoted config
    # argument breaks on paths containing spaces and kills every OCR attempt.
    td = tessdata_dir()
    if td:
        os.environ["TESSDATA_PREFIX"] = td
    psm_modes = [3, 4, 6, 11, 12]

    def data_confidence(image, lang, psm) -> tuple[float, Optional[dict]]:
        """Tesseract's own per-word confidence (real text ~70-95, glyph noise
        ~20-45) plus word bounding boxes as page-relative fractions, used for
        search-hit highlighting in the page viewer."""
        try:
            data = pytesseract.image_to_data(
                image, lang=lang, config=f"--psm {psm}", output_type=pytesseract.Output.DICT
            )
            iw, ih = (getattr(image, "width", 0) or 1), (getattr(image, "height", 0) or 1)
            confs = []
            words = []
            n = len(data.get("text", []))
            for i in range(n):
                w = str(data["text"][i]).strip()
                try:
                    c = int(float(data["conf"][i]))
                except (TypeError, ValueError):
                    continue
                if c < 0 or not w:
                    continue
                confs.append(c)
                words.append({
                    "t": w,
                    "x0": round(data["left"][i] / iw, 4),
                    "y0": round(data["top"][i] / ih, 4),
                    "x1": round((data["left"][i] + data["width"][i]) / iw, 4),
                    "y1": round((data["top"][i] + data["height"][i]) / ih, 4),
                })
            conf = (sum(confs) / len(confs)) if confs else 0.0
            return conf, ({"words": words} if words else None)
        except Exception:
            return -1.0, None

    def run_grid(image):
        top = {"text": "", "score": 0, "engine": "", "conf": -1.0, "words": None, "_probe": None}
        for variant_name, variant in ocr_image_variants(image):
            for lang in attempts:
                for psm in psm_modes:
                    try:
                        txt = pytesseract.image_to_string(variant, lang=lang, config=f"--psm {psm}")
                        score = ocr_text_score(txt)
                        if score > top["score"]:
                            top = {"text": txt, "score": score, "engine": f"tesseract:{lang}:psm{psm}:{variant_name}", "conf": -1.0, "words": None, "_probe": (variant, lang, psm)}
                            # A clean, high-confidence read ends the search early;
                            # otherwise every page runs the whole ~100-pass grid.
                            if ocr_quality(txt)[0] == "good":
                                conf, words = data_confidence(variant, lang, psm)
                                top["conf"] = conf
                                top["words"] = words
                                if conf < 0 or conf >= 60:
                                    top.pop("_probe", None)
                                    return top
                    except Exception as exc:
                        errors.append(f"{lang}/psm{psm}/{variant_name}: {exc}")
        # No early exit: measure the winner's confidence once anyway, so
        # fallback decisions (vision models, rotation) know how sure the
        # engine actually was - "medium"-quality garbage must not pass silently.
        probe = top.pop("_probe", None)
        if top["text"].strip() and top["conf"] < 0 and probe:
            conf, words = data_confidence(*probe)
            top["conf"] = conf
            top["words"] = words
        return top

    best = run_grid(img)

    # Sideways or upside-down scans (common with phone photos): if the straight
    # read failed or came back as low-confidence noise, ask Tesseract's
    # orientation detector and retry once on the corrected rotation.
    needs_rescue = ocr_quality(best["text"])[0] in ("empty", "poor") or (0 <= best["conf"] < 50)
    if needs_rescue and "osd" in langs:
        try:
            osd = pytesseract.image_to_osd(img)
            m = re.search(r"Rotate:\s*(\d+)", osd)
            angle = int(m.group(1)) if m else 0
            if angle:
                rotated = img.rotate(-angle, expand=True)
                r = run_grid(rotated)
                better_conf = r["conf"] >= 0 and r["conf"] > max(best["conf"], 0.0) + 10
                if better_conf or (r["conf"] < 0 and r["score"] > best["score"]):
                    r["engine"] += f":rot{angle}"
                    # Boxes were measured on the rotated raster and would land in
                    # the wrong place on the original page image.
                    r["words"] = None
                    best = r
        except Exception as exc:
            errors.append(f"osd: {exc}")

    if best["text"].strip():
        engine = f"{best['engine']}:conf{int(best['conf'])}" if best["conf"] >= 0 else best["engine"]
        return best["text"], engine, best.get("words")
    return "", "; ".join(errors) or "tesseract-no-text", None


def ocr_image_openai(path: Path) -> tuple[str, str]:
    if not (os.environ.get("OPENAI_API_KEY") and OpenAI):
        return "", "openai-not-configured"
    try:
        client = OpenAI()
        mime = "image/png" if path.suffix.lower() == ".png" else "image/jpeg"
        data = base64.b64encode(path.read_bytes()).decode("ascii")
        prompt = (
            "Extract all visible text from this document image. Support Arabic and English. "
            "Preserve line breaks, numbers, dates, totals, and names. Return only extracted text."
        )
        resp = client.chat.completions.create(
            model=os.environ.get("DOCWISE_VISION_MODEL", "gpt-4o-mini"),
            messages=[{"role": "user", "content": [
                {"type": "text", "text": prompt},
                {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{data}"}},
            ]}],
            temperature=0,
        )
        return resp.choices[0].message.content or "", "openai-vision"
    except Exception as exc:
        return "", f"openai-error: {exc}"


def ocr_image_azure(path: Path) -> tuple[str, str]:
    """Azure Document Intelligence prebuilt-read: strong Arabic print and
    handwriting. Used as a low-confidence fallback only when the user sets
    AZURE_DI_ENDPOINT and AZURE_DI_KEY (paid per page)."""
    endpoint = (os.environ.get("AZURE_DI_ENDPOINT") or "").rstrip("/")
    key = os.environ.get("AZURE_DI_KEY") or ""
    if not (endpoint and key):
        return "", "azure-not-configured"
    try:
        import urllib.request

        url = f"{endpoint}/documentintelligence/documentModels/prebuilt-read:analyze?api-version=2024-11-30"
        req = urllib.request.Request(
            url,
            data=path.read_bytes(),
            headers={"Ocp-Apim-Subscription-Key": key, "Content-Type": "application/octet-stream"},
            method="POST",
        )
        with urllib.request.urlopen(req, timeout=60) as r:
            op_url = r.headers.get("Operation-Location")
        if not op_url:
            return "", "azure-error: no operation location"
        for _ in range(60):
            time.sleep(1.5)
            poll = urllib.request.Request(op_url, headers={"Ocp-Apim-Subscription-Key": key})
            with urllib.request.urlopen(poll, timeout=30) as r:
                out = json.loads(r.read().decode())
            state = out.get("status")
            if state == "succeeded":
                content = ((out.get("analyzeResult") or {}).get("content") or "").strip()
                return content, "azure-document-intelligence"
            if state in ("failed", "canceled"):
                return "", f"azure-{state}"
        return "", "azure-timeout"
    except Exception as exc:
        return "", f"azure-error: {str(exc)[:100]}"


_OLLAMA_CACHE = {"ts": 0.0, "status": None}


def ollama_status() -> dict:
    """Detect a local Ollama server and its best vision model (cached 60s)."""
    if time.time() - _OLLAMA_CACHE["ts"] < 60 and _OLLAMA_CACHE["status"] is not None:
        return _OLLAMA_CACHE["status"]
    base = os.environ.get("DOCWISE_OLLAMA_URL", "http://127.0.0.1:11434")
    status = {"available": False, "model": None}
    if os.environ.get("DOCWISE_OLLAMA", "auto").lower() not in ("0", "off", "false"):
        try:
            import urllib.request
            with urllib.request.urlopen(f"{base}/api/tags", timeout=2) as r:
                tags = json.loads(r.read().decode())
            names = [m.get("name", "") for m in tags.get("models", [])]
            forced = os.environ.get("DOCWISE_OLLAMA_MODEL", "")
            if forced:
                model = forced
            else:
                vision_pref = ("qwen2.5vl", "qwen3-vl", "llama3.2-vision", "minicpm-v", "llava", "moondream", "gemma3")
                model = next((n for n in names if any(n.lower().startswith(p) for p in vision_pref)), None)
            status = {"available": bool(model), "model": model}
        except Exception:
            pass
    _OLLAMA_CACHE.update({"ts": time.time(), "status": status})
    return status


def ollama_chat_model() -> Optional[str]:
    """Pick a local text model for answering: forced via env, else best match."""
    forced = os.environ.get("DOCWISE_OLLAMA_CHAT_MODEL", "")
    if forced:
        return forced
    base = os.environ.get("DOCWISE_OLLAMA_URL", "http://127.0.0.1:11434")
    if os.environ.get("DOCWISE_OLLAMA", "auto").lower() in ("0", "off", "false"):
        return None
    try:
        import urllib.request
        with urllib.request.urlopen(f"{base}/api/tags", timeout=2) as r:
            tags = json.loads(r.read().decode())
        models = [(m.get("name", ""), int(m.get("size") or 0)) for m in tags.get("models", [])]
        bad = ("coder", "embed", "bge", "nomic", "cloud", "vision", "moondream", "llava")
        eligible = [(n, s) for n, s in models if n and not any(b in n.lower() for b in bad)]
        for pref in ("qwen3.5", "qwen3", "qwen2.5", "gemma3", "gemma4", "llama3", "mistral", "hermes", "phi"):
            sized = [(n, s) for n, s in eligible if n.lower().startswith(pref) and s > 2_000_000_000]
            if sized:
                # Largest of the preferred family under ~11GB keeps quality high
                # and load times sane on normal machines.
                capped = [x for x in sized if x[1] <= 11_000_000_000] or sized
                return max(capped, key=lambda x: x[1])[0]
        return eligible[0][0] if eligible else None
    except Exception:
        return None


def ollama_answer(question: str, context: str, lang_hint: str) -> Optional[str]:
    """Compose a grounded answer with a local Ollama text model."""
    model = ollama_chat_model()
    if not model:
        return None
    base = os.environ.get("DOCWISE_OLLAMA_URL", "http://127.0.0.1:11434")
    prompt = (
        "You are a strict document archive assistant. Use ONLY the provided sources. "
        "Cite sources like [Source 1]. If the sources do not clearly support an answer, "
        f"say you could not find it in the indexed documents. Prefer exact amounts, dates and names. {lang_hint}\n\n"
        f"Question: {question}\n\nSources:\n{context}\n\n"
        f"IMPORTANT: {lang_hint}\nAnswer:"
    )
    try:
        import urllib.request
        payload = json.dumps({
            "model": model,
            "prompt": prompt,
            "stream": False,
            "think": False,
            "options": {"temperature": 0.2, "num_ctx": 8192},
        }).encode()
        req = urllib.request.Request(f"{base}/api/generate", data=payload, headers={"Content-Type": "application/json"})
        with urllib.request.urlopen(req, timeout=int(os.environ.get("DOCWISE_OLLAMA_TIMEOUT", "240"))) as r:
            out = json.loads(r.read().decode())
        answer = (out.get("response") or "").strip()
        answer = re.sub(r"<think>.*?</think>", "", answer, flags=re.S).strip()
        return answer or None
    except Exception:
        return None


def ocr_image_ollama(path: Path) -> tuple[str, str]:
    """OCR via a local Ollama vision model: free, private, no API key."""
    st = ollama_status()
    if not st["available"]:
        return "", "ollama-unavailable"
    base = os.environ.get("DOCWISE_OLLAMA_URL", "http://127.0.0.1:11434")
    try:
        import urllib.request
        payload = json.dumps({
            "model": st["model"],
            "prompt": (
                "Extract all visible text from this document image exactly as written. "
                "Support Arabic and English. Preserve line breaks, numbers, dates and totals. "
                "Return only the extracted text, nothing else."
            ),
            "images": [base64.b64encode(path.read_bytes()).decode("ascii")],
            "stream": False,
            "options": {"temperature": 0},
        }).encode()
        req = urllib.request.Request(f"{base}/api/generate", data=payload, headers={"Content-Type": "application/json"})
        with urllib.request.urlopen(req, timeout=int(os.environ.get("DOCWISE_OLLAMA_TIMEOUT", "240"))) as r:
            out = json.loads(r.read().decode())
        return (out.get("response") or "").strip(), f"ollama:{st['model']}"
    except Exception as exc:
        return "", f"ollama-error: {str(exc)[:80]}"


def ocr_image(path: Path) -> tuple[str, str, Optional[dict]]:
    # DOCWISE_OCR forces one engine (tesseract/paddle/ollama/azure/openai);
    # a forced engine that produces nothing falls through to the auto chain.
    engine = os.environ.get("DOCWISE_OCR", "auto").lower()
    if engine in ("openai", "vision"):
        txt, used = ocr_image_openai(path)
        if txt.strip():
            return txt, used, None
    if engine == "ollama":
        txt, used = ocr_image_ollama(path)
        if txt.strip():
            return txt, used, None
    if engine == "azure":
        txt, used = ocr_image_azure(path)
        if txt.strip():
            return txt, used, None
    if engine == "tesseract":
        return ocr_image_tesseract(path)

    txt, used, words = ocr_image_tesseract(path)
    m = re.search(r":conf(\d+)", used or "")
    t_conf = float(m.group(1)) if m else -1.0

    if txt.strip() and not should_try_vision_fallback(txt, used):
        return txt, used, words

    # Hard page (IDs, stamps, handwriting, busy backgrounds). Fallback order:
    # local vision model (free, private) -> Azure Document Intelligence (cheap,
    # OCR-specialized, only if configured) -> OpenAI Vision. When the engine
    # itself was unsure, a clean fallback reading beats a higher glyph count -
    # garbage salad outscores correct text on raw character counting.
    low_conf = not txt.strip() or (0 <= t_conf < 65) or ocr_quality(txt)[0] in ("empty", "poor")

    def fallback_wins(candidate: str) -> bool:
        if not candidate.strip():
            return False
        if ocr_text_score(candidate) >= ocr_text_score(txt):
            return True
        return low_conf and len(candidate.strip()) >= 20 and not text_is_garbled(candidate)

    txt3, used3 = ocr_image_ollama(path)
    if fallback_wins(txt3):
        return txt3, f"{used3}:fallback-from-{used}", None
    txt4, used4 = ocr_image_azure(path)
    if fallback_wins(txt4):
        return txt4, f"{used4}:fallback-from-{used}", None
    txt2, used2 = ocr_image_openai(path)
    if fallback_wins(txt2):
        return txt2, f"{used2}:fallback-from-{used}", None
    if txt.strip():
        return txt, used, words
    return "", f"{used}; {used3}; {used4}; {used2}", None


OCR_STATE = {"running": False, "file": "", "pages_done": 0, "pages_total": 0}


def extract_text_from_pdf(path: Path) -> tuple[list[dict], str, Optional[str]]:
    pages = []
    engines = set()
    errors = []
    if not fitz:
        detail = FITZ_IMPORT_ERROR or "PyMuPDF is not installed"
        return [], "none", (
            f"PyMuPDF cannot load: {detail}. On Windows this is usually fixed by "
            "installing the Microsoft Visual C++ runtime "
            "(https://aka.ms/vs/17/release/vc_redist.x64.exe) or re-running setup.ps1."
        )
    try:
        doc = fitz.open(path)
        max_ocr_pages = int(os.environ.get("DOCWISE_MAX_OCR_PDF_PAGES", "20"))
        # Pass 1: pull embedded text and render scan pages to temp images.
        # Rendering stays sequential (fitz documents are not thread-safe);
        # the slow part - Tesseract - runs in parallel below.
        ocr_jobs = []
        for i, page in enumerate(doc, start=1):
            text = page.get_text("text") or ""
            entry = {"page": i, "text": text, "words": None, "engine": "pdf-text"}
            pages.append(entry)
            if should_ocr_pdf_page(text) and i <= max_ocr_pages:
                try:
                    # 3.5x ~= 250 dpi: measurably better Arabic diacritics and
                    # small print than the old 2.5x, at no real speed cost.
                    pix = page.get_pixmap(matrix=fitz.Matrix(3.5, 3.5), alpha=False)
                    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                        tmp_path = Path(tmp.name)
                    pix.save(str(tmp_path))
                    ocr_jobs.append((entry, tmp_path))
                except Exception as exc:
                    errors.append(f"page {i} render failed: {exc}")

        if ocr_jobs:
            from concurrent.futures import ThreadPoolExecutor

            workers = max(1, int(os.environ.get("DOCWISE_OCR_WORKERS", str(min(4, (os.cpu_count() or 2))))))
            OCR_STATE.update({"running": True, "file": path.name, "pages_done": 0, "pages_total": len(ocr_jobs)})

            def ocr_one(job):
                entry, tmp_path = job
                try:
                    ocr_text, ocr_engine, ocr_words = ocr_image(tmp_path)
                    if ocr_text.strip() and (not text_is_garbled(ocr_text) or len(ocr_text.strip()) > len(entry["text"].strip())):
                        entry["text"] = ocr_text
                        entry["engine"] = f"pdf-render+{ocr_engine}"
                        entry["words"] = ocr_words
                    else:
                        errors.append(f"page {entry['page']}: {ocr_engine}")
                except Exception as exc:
                    errors.append(f"page {entry['page']} OCR failed: {exc}")
                finally:
                    tmp_path.unlink(missing_ok=True)
                    OCR_STATE["pages_done"] += 1

            try:
                with ThreadPoolExecutor(max_workers=workers) as pool:
                    list(pool.map(ocr_one, ocr_jobs))
            finally:
                OCR_STATE.update({"running": False, "file": ""})

        for entry in pages:
            engines.add(entry.pop("engine", "pdf-text"))
        return pages, ",".join(sorted(engines)) or "pdf", "; ".join(errors) if errors else None
    except Exception as exc:
        return [], "pdf-error", str(exc)


def strip_rtf(text: str) -> str:
    text = re.sub(r"\\'[0-9a-fA-F]{2}", " ", text)
    text = re.sub(r"\\[a-zA-Z]+-?\d* ?", " ", text)
    text = text.replace("{", " ").replace("}", " ")
    return re.sub(r"\s+", " ", text).strip()


def extract_word_text(path: Path) -> tuple[list[dict], str, Optional[str]]:
    if not Document:
        return [], "docx-error", "python-docx is not installed"
    try:
        doc = Document(str(path))
        parts = []
        for p in doc.paragraphs:
            if p.text and p.text.strip():
                parts.append(p.text.strip())
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text and cell.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))
        return [{"page": 1, "text": "\n".join(parts)}], "docx-text", None
    except Exception as exc:
        return [], "docx-error", str(exc)


def extract_excel_text(path: Path) -> tuple[list[dict], str, Optional[str]]:
    if not load_workbook:
        return [], "xlsx-error", "openpyxl is not installed"
    try:
        wb = load_workbook(str(path), read_only=True, data_only=True)
        pages = []
        for i, ws in enumerate(wb.worksheets, start=1):
            lines = [f"Sheet: {ws.title}"]
            for row in ws.iter_rows(values_only=True):
                values = [str(v).strip() for v in row if v is not None and str(v).strip()]
                if values:
                    lines.append(" | ".join(values))
            pages.append({"page": i, "text": "\n".join(lines)})
        wb.close()
        return pages, "xlsx-text", None
    except Exception as exc:
        return [], "xlsx-error", str(exc)


def extract_powerpoint_text(path: Path) -> tuple[list[dict], str, Optional[str]]:
    if not Presentation:
        return [], "pptx-error", "python-pptx is not installed"
    try:
        prs = Presentation(str(path))
        pages = []
        for i, slide in enumerate(prs.slides, start=1):
            parts = [f"Slide {i}"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and shape.text.strip():
                    parts.append(shape.text.strip())
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text and cell.text.strip()]
                        if cells:
                            parts.append(" | ".join(cells))
            pages.append({"page": i, "text": "\n".join(parts)})
        return pages, "pptx-text", None
    except Exception as exc:
        return [], "pptx-error", str(exc)


def extract_legacy_office_text(path: Path) -> tuple[list[dict], str, Optional[str]]:
    """Best-effort extraction for old .doc/.xls/.ppt using installed Microsoft Office COM."""
    try:
        import win32com.client  # type: ignore
        ext = path.suffix.lower()
        if ext == ".doc":
            word = win32com.client.DispatchEx("Word.Application")
            word.Visible = False
            doc = word.Documents.Open(str(path.resolve()), ReadOnly=True)
            text = doc.Content.Text
            doc.Close(False)
            word.Quit()
            return [{"page": 1, "text": text}], "word-com-text", None
        if ext == ".xls":
            excel = win32com.client.DispatchEx("Excel.Application")
            excel.Visible = False
            wb = excel.Workbooks.Open(str(path.resolve()), ReadOnly=True)
            pages = []
            for i, ws in enumerate(wb.Worksheets, start=1):
                used = ws.UsedRange.Value
                lines = [f"Sheet: {ws.Name}"]
                if isinstance(used, tuple):
                    for row in used:
                        if not isinstance(row, tuple):
                            row = (row,)
                        values = [str(v).strip() for v in row if v is not None and str(v).strip()]
                        if values:
                            lines.append(" | ".join(values))
                elif used is not None:
                    lines.append(str(used))
                pages.append({"page": i, "text": "\n".join(lines)})
            wb.Close(False)
            excel.Quit()
            return pages, "excel-com-text", None
        if ext == ".ppt":
            ppt = win32com.client.DispatchEx("PowerPoint.Application")
            pres = ppt.Presentations.Open(str(path.resolve()), WithWindow=False)
            pages = []
            for i, slide in enumerate(pres.Slides, start=1):
                parts = [f"Slide {i}"]
                for shape in slide.Shapes:
                    try:
                        if shape.HasTextFrame and shape.TextFrame.HasText:
                            parts.append(shape.TextFrame.TextRange.Text)
                    except Exception:
                        pass
                pages.append({"page": i, "text": "\n".join(parts)})
            pres.Close()
            ppt.Quit()
            return pages, "powerpoint-com-text", None
    except Exception as exc:
        return [], "legacy-office-error", f"Legacy Office extraction failed. Install Microsoft Office or convert to docx/xlsx/pptx. Details: {exc}"
    return [], "legacy-office-error", "Unsupported legacy Office file"


def extract_text(path: Path) -> tuple[list[dict], str, Optional[str]]:
    ext = path.suffix.lower()
    if ext in TEXT_EXTS:
        for enc in ("utf-8", "utf-8-sig", "cp1256", "latin-1"):
            try:
                return [{"page": 1, "text": path.read_text(encoding=enc, errors="ignore")}], f"text:{enc}", None
            except Exception:
                pass
        return [], "text-error", "Could not decode text file"
    if ext == ".rtf":
        for enc in ("utf-8", "utf-8-sig", "cp1256", "latin-1"):
            try:
                return [{"page": 1, "text": strip_rtf(path.read_text(encoding=enc, errors="ignore"))}], f"rtf:{enc}", None
            except Exception:
                pass
        return [], "rtf-error", "Could not decode RTF file"
    if ext == ".pdf":
        return extract_text_from_pdf(path)
    if ext in IMAGE_EXTS:
        text, engine, words = ocr_image(path)
        return [{"page": 1, "text": text, "words": words}], engine, None if text.strip() else engine
    if ext in WORD_EXTS:
        return extract_word_text(path)
    if ext in EXCEL_EXTS:
        return extract_excel_text(path)
    if ext in POWERPOINT_EXTS:
        return extract_powerpoint_text(path)
    if ext in LEGACY_OFFICE_EXTS:
        return extract_legacy_office_text(path)
    return [], "unsupported", "Unsupported file type"


def guess_metadata(path: Path, text: str) -> dict:
    lowered = normalize_arabic(text + " " + path.name)
    raw = text or ""
    lang = "Arabic + English" if has_arabic(raw) and re.search(r"[A-Za-z]", raw) else "Arabic" if has_arabic(raw) else "English" if re.search(r"[A-Za-z]", raw) else "unknown"

    # Order matters: specific document types first. The invoice rule contains
    # generic money words (total, riyal...) that appear inside bank statements,
    # contracts and receipts - it must run after every specific type.
    rules = [
        # Forms above id: application forms ask FOR an Emirates ID number,
        # but real ID documents never say "application form".
        ("form", ["application form", "نموذج طلب", "استماره", "نموذج تسجيل"]),
        ("id", ["passport", "emirates id", "national id", "identity card", "جواز", "هويه", "بطاقه", "اقامه"]),
        # "bank"/"بنك" alone are too generic: receipts and invoices constantly
        # say "bank transfer". Statement-specific markers only.
        ("bank", ["bank statement", "account statement", "iban", "swift", "كشف حساب", "ايبان", "الرصيد الافتتاحي", "closing balance"]),
        ("certificate", ["certificate", "شهاده", "degree", "diploma", "اتم بنجاح", "دوره تدريبيه", "completion", "يشهد", "تشهد"]),
        # Receipts before contracts: rent receipts mention the lease, but
        # contracts never carry receipt vouchers.
        ("receipt", ["receipt", "ايصال", "سند قبض", "payment received", "مدفوع", "استلمنا"]),
        # Contracts before medical/legal: employment contracts routinely
        # mention medical insurance and legal clauses.
        ("contract", ["contract", "agreement", "عقد", "اتفاقيه", "lease", "ايجار", "nda", "tenancy", "ejari", "ايجاري"]),
        ("medical", ["medical", "hospital", "clinic", "doctor", "patient", "طبي", "مستشفي", "عياده", "مريض", "وصفه", "prescription", "lab report", "تقرير طبي"]),
        ("legal", ["court", "legal", "law", "محكمه", "قانون", "دعوي"]),
        ("invoice", ["invoice", "فاتوره", "bill", "amount due", "total", "ضريبه", "vat", "aed", "درهم", "ريال", "sar", "الاجمالي"]),
        ("news", ["news", "article", "مقال", "خبر", "اخبار", "كشفت", "اعلنت", "تقنيه"]),
    ]
    doc_type = "general"
    for typ, kws in rules:
        if any(contains_keyword(lowered, k) for k in kws):
            doc_type = typ
            break

    raw_d = normalize_digits(raw)
    date_guess = None
    date_patterns = [
        r"\b(20\d{2}[-/.]\d{1,2}[-/.]\d{1,2})\b",
        r"\b(\d{1,2}[-/.]\d{1,2}[-/.]20\d{2})\b",
        r"\b(\d{1,2}\s+(?:jan|feb|mar|apr|may|jun|jul|aug|sep|oct|nov|dec)[a-z]*\s+20\d{2})\b",
        r"\b(14\d{2}[-/.]\d{1,2}[-/.]\d{1,2})\b",
    ]
    for pat in date_patterns:
        m = re.search(pat, raw_d, re.I)
        if m:
            date_guess = m.group(1)
            break

    amount = None
    amount_patterns = [rf"{CURRENCY_TOKEN}\s*({AMOUNT_NUM})", rf"({AMOUNT_NUM})\s*{CURRENCY_TOKEN}"]
    for pat in amount_patterns:
        matches = re.findall(pat, raw_d, re.I)
        if matches:
            amount = matches[-1]
            break

    company = None
    lines = [ln.strip() for ln in raw.splitlines() if ln.strip()]
    for ln in lines[:12]:
        if 2 <= len(ln) <= 80 and not text_is_garbled(ln) and not re.search(r"^[\d\W_]+$", ln):
            company = ln[:80]
            break

    title = path.stem.replace("_", " ").replace("-", " ").strip().title()
    tags = [doc_type, lang.lower().split()[0]]
    if amount:
        tags.append("amount")
    if date_guess:
        tags.append("date")

    summary = summarize_text(text, doc_type, lang, date_guess, amount, company)
    return {
        "title": title,
        "doc_type": doc_type,
        "language": lang,
        "date_guess": date_guess,
        "company": company,
        "amount": amount,
        "tags": sorted(set(tags)),
        "summary": summary,
    }


def summarize_text(text: str, doc_type: str, lang: str, date_guess: Optional[str], amount: Optional[str], company: Optional[str]) -> str:
    clean = re.sub(r"\s+", " ", text or "").strip()
    if not clean:
        return "No text extracted yet. Install Tesseract with Arabic language data or set OPENAI_API_KEY for image OCR."
    prefix = f"{doc_type.title()} document"
    if company:
        prefix += f" from {company}"
    bits = []
    if date_guess:
        bits.append(f"date: {date_guess}")
    if amount:
        bits.append(f"amount: {amount}")
    if lang != "unknown":
        bits.append(f"language: {lang}")
    meta = f" ({', '.join(bits)})" if bits else ""
    return f"{prefix}{meta}. {clean[:420]}{'...' if len(clean) > 420 else ''}"


def extract_structured_fields(text: str, meta: dict) -> dict:
    """Document-type fields used before/alongside RAG answers."""
    raw = text or ""
    fields = {
        "doc_type": meta.get("doc_type"),
        "title": meta.get("title"),
        "company": meta.get("company"),
        "date": meta.get("date_guess"),
        "amount": meta.get("amount"),
    }
    # Common fields (Arabic-Indic digits normalized; ة/ه spelling variants covered)
    raw_d = normalize_digits(raw)
    m = re.search(r"(?:invoice\s*(?:no|number|#)|رقم\s*الفاتور[ةه]|فاتور[ةه]\s*رقم)\s*[:#-]?\s*([A-Z0-9\-/]+)", raw_d, re.I)
    if m:
        fields["invoice_number"] = m.group(1)
    m = re.search(r"(?:due\s*date|expiry\s*date|تاريخ\s*(?:الاستحقاق|الانتهاء|انتهاء))\s*[:#-]?\s*([0-9A-Za-z\-/\.\s]+)", raw_d, re.I)
    if m:
        fields["due_or_expiry_date"] = m.group(1).strip()[:60]
    amounts = re.findall(rf"{CURRENCY_TOKEN}\s*({AMOUNT_NUM})|({AMOUNT_NUM})\s*{CURRENCY_TOKEN}", raw_d, re.I)
    clean_amounts = [a or b for a, b in amounts if (a or b)]
    if clean_amounts:
        fields["amounts_found"] = clean_amounts[-8:]
        fields["amount"] = fields.get("amount") or clean_amounts[-1]
    cur = re.search(CURRENCY_TOKEN, raw_d, re.I)
    if cur:
        fields["currency"] = cur.group(0)
    # Prefer the amount on a "total" line as THE total.
    total_kw = re.compile(r"(?:grand\s*total|total\s*(?:due|amount)?|الاجمالي|الإجمالي|المجموع|المبلغ\s*الاجمالي|المبلغ\s*المستحق|صافي)", re.I)
    for line in raw_d.splitlines():
        if total_kw.search(line):
            lm = re.search(rf"({AMOUNT_NUM})\s*{CURRENCY_TOKEN}|{CURRENCY_TOKEN}\s*({AMOUNT_NUM})|({AMOUNT_NUM})\s*$", line.strip(), re.I)
            if lm:
                fields["total"] = next((g for g in lm.groups() if g), None)
                break
    m = re.search(r"(?:vat\s*(?:no|number|reg(?:istration)?)?|tax\s*(?:no|number)|الرقم\s*الضريبي|رقم\s*ضريبي)\s*[:#.]?\s*(\d{9,15})", raw_d, re.I)
    if m:
        fields["vat_number"] = m.group(1)
    m = re.search(r"(?:iban|ايبان|آيبان)\s*[:#-]?\s*([A-Z]{2}\s?[0-9]{2}[0-9A-Z\s]{10,32})", raw_d, re.I)
    if m:
        fields["iban"] = re.sub(r"\s+", "", m.group(1))[:34]
    if meta.get("doc_type") == "contract":
        rent = re.search(rf"(?:rent|annual rent|الايجار|قيمة الايجار|الأجر[ةه]|قيم[ةه]\s*الإيجار)\s*[:#-]?\s*({AMOUNT_NUM}\s*{CURRENCY_TOKEN}?)", raw_d, re.I)
        if rent:
            fields["rent_amount"] = rent.group(1)
    if meta.get("doc_type") in ("id", "legal", "contract"):
        names = re.findall(r"(?:name|الاسم)\s*[:#-]?\s*([^\n]{3,80})", raw, re.I)
        if names:
            fields["names_found"] = [n.strip() for n in names[:6]]
    # Optional schema extraction with OpenAI for better invoices/contracts/IDs.
    if os.environ.get("OPENAI_API_KEY") and OpenAI and os.environ.get("DOCWISE_FIELD_AI", "1") == "1" and raw.strip():
        try:
            client = OpenAI()
            schema_prompt = (
                "Extract structured fields from this document as compact JSON only. "
                "Support Arabic/English. Include only fields clearly present. "
                "Useful keys: vendor, customer, invoice_number, date, due_date, total, currency, tax, parties, start_date, end_date, id_number, expiry_date, category."
            )
            resp = client.chat.completions.create(
                model=os.environ.get("DOCWISE_CHAT_MODEL", "gpt-4o-mini"),
                messages=[{"role": "user", "content": f"{schema_prompt}\n\nDocument type: {meta.get('doc_type')}\nText:\n{raw[:6000]}"}],
                temperature=0,
            )
            content = resp.choices[0].message.content or "{}"
            content = re.sub(r"^```json|```$", "", content.strip(), flags=re.I).strip()
            ai_fields = json.loads(content)
            if isinstance(ai_fields, dict):
                fields["ai"] = ai_fields
        except Exception as exc:
            fields["ai_error"] = str(exc)[:160]
    return {k: v for k, v in fields.items() if v not in (None, "", [], {})}


def chunk_pages(pages: list[dict], max_tokens: int = 700, overlap_tokens: int = 90) -> list[dict]:
    """Structure-aware-ish chunking: paragraphs/lines first, token window fallback, page citations preserved."""
    chunks = []
    for page in pages:
        raw = (page.get("text") or "").strip()
        if not raw:
            continue
        # Keep table rows/slide lines together when possible; don't flatten everything immediately.
        blocks = [b.strip() for b in re.split(r"\n\s*\n|(?<=\.)\s+(?=[A-Z\u0600-\u06FF])", raw) if b.strip()]
        if len(blocks) <= 1:
            blocks = [b.strip() for b in raw.splitlines() if b.strip()] or [raw]
        idx = 0
        current = []
        current_tokens = 0
        for block in blocks:
            bt = approx_token_count(block)
            if current and current_tokens + bt > max_tokens:
                text = "\n".join(current).strip()
                chunks.append({"page": page.get("page", 1), "chunk_index": idx, "text": text, "token_count": approx_token_count(text)})
                idx += 1
                # Lightweight overlap by carrying trailing lines.
                carry = []
                carry_tokens = 0
                for prev in reversed(current):
                    pt = approx_token_count(prev)
                    if carry_tokens + pt > overlap_tokens:
                        break
                    carry.insert(0, prev)
                    carry_tokens += pt
                current = carry
                current_tokens = carry_tokens
            if bt > max_tokens:
                words = block.split()
                step = max(1, int(max_tokens - overlap_tokens))
                for start in range(0, len(words), step):
                    part = " ".join(words[start:start + max_tokens]).strip()
                    if part:
                        chunks.append({"page": page.get("page", 1), "chunk_index": idx, "text": part, "token_count": approx_token_count(part)})
                        idx += 1
                current = []
                current_tokens = 0
            else:
                current.append(block)
                current_tokens += bt
        if current:
            text = "\n".join(current).strip()
            chunks.append({"page": page.get("page", 1), "chunk_index": idx, "text": text, "token_count": approx_token_count(text)})
    return chunks


def index_file(path: Path, source_type: str = "folder", force: bool = False) -> dict:
    path = path.resolve()
    if path.suffix.lower() not in SUPPORTED or not path.is_file():
        return {"status": "skipped", "path": str(path), "reason": "unsupported"}
    stat = path.stat()
    sha = file_hash(path)

    with DB_LOCK, db() as conn:
        existing = conn.execute("SELECT id, sha256, mtime FROM documents WHERE path=?", (str(path),)).fetchone()
        if existing and not force and existing["sha256"] == sha and abs((existing["mtime"] or 0) - stat.st_mtime) < 0.001:
            return {"status": "unchanged", "id": existing["id"], "path": str(path)}
        # Same content already indexed under another path (copied file, second
        # folder, re-upload): skip instead of indexing a duplicate. Only a
        # healthy existing copy blocks - a failed/empty one never should.
        if not existing and not force and os.environ.get("DOCWISE_SKIP_DUPLICATES", "1") == "1":
            dup = conn.execute(
                "SELECT id, path FROM documents WHERE sha256=? AND status='indexed' LIMIT 1", (sha,)
            ).fetchone()
            if dup and Path(dup["path"]).exists():
                return {"status": "duplicate", "id": dup["id"], "path": str(path), "duplicate_of": dup["path"]}

    pages, engine, err = extract_text(path)
    full_text = "\n\n".join([p.get("text") or "" for p in pages]).strip()
    meta = guess_metadata(path, full_text)
    fields = extract_structured_fields(full_text, meta)
    quality, quality_score = ocr_quality(full_text)
    normalized = normalize_arabic(full_text + " " + meta["title"] + " " + " ".join(meta["tags"]) + " " + json.dumps(fields, ensure_ascii=False))
    status = "indexed" if full_text else "needs_ocr"
    if err and not full_text:
        status = "error" if "Unsupported" in err else "needs_ocr"

    with DB_LOCK, db() as conn:
        cur = conn.execute(
            """
            INSERT INTO documents(path, original_name, source_type, file_ext, size, mtime, sha256, title, doc_type, language,
                date_guess, company, amount, tags, summary, fields, text, normalized_text, status, ocr_engine, ocr_quality, ocr_score, error, created_at, updated_at)
            VALUES(?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?)
            ON CONFLICT(path) DO UPDATE SET
                original_name=excluded.original_name, source_type=excluded.source_type, file_ext=excluded.file_ext,
                size=excluded.size, mtime=excluded.mtime, sha256=excluded.sha256, title=excluded.title,
                doc_type=excluded.doc_type, language=excluded.language, date_guess=excluded.date_guess,
                company=excluded.company, amount=excluded.amount, tags=excluded.tags, summary=excluded.summary, fields=excluded.fields,
                text=excluded.text, normalized_text=excluded.normalized_text, status=excluded.status,
                ocr_engine=excluded.ocr_engine, ocr_quality=excluded.ocr_quality, ocr_score=excluded.ocr_score, error=excluded.error, updated_at=excluded.updated_at
            """,
            (
                str(path), path.name, source_type, path.suffix.lower(), stat.st_size, stat.st_mtime, sha,
                meta["title"], meta["doc_type"], meta["language"], meta["date_guess"], meta["company"], meta["amount"],
                json.dumps(meta["tags"], ensure_ascii=False), meta["summary"], json.dumps(fields, ensure_ascii=False),
                full_text, normalized, status, engine, quality, quality_score, err, now_iso(), now_iso(),
            ),
        )
        doc_id_row = conn.execute("SELECT id FROM documents WHERE path=?", (str(path),)).fetchone()
        doc_id = doc_id_row["id"]
        delete_chunk_indexes_for_doc(conn, doc_id)
        conn.execute("DELETE FROM chunks WHERE document_id=?", (doc_id,))
        conn.execute("DELETE FROM page_words WHERE document_id=?", (doc_id,))
        for p in pages:
            if p.get("words") and p["words"].get("words"):
                conn.execute(
                    "INSERT OR REPLACE INTO page_words(document_id, page, words) VALUES(?,?,?)",
                    (doc_id, p.get("page", 1), json.dumps(p["words"]["words"], ensure_ascii=False)),
                )
        for ch in chunk_pages(pages):
            norm_ch = normalize_arabic(ch["text"])
            cur = conn.execute(
                "INSERT INTO chunks(document_id, page, chunk_index, text, normalized_text, token_count) VALUES(?,?,?,?,?,?)",
                (doc_id, ch["page"], ch["chunk_index"], ch["text"], norm_ch, ch.get("token_count") or approx_token_count(ch["text"])),
            )
            add_chunk_indexes(conn, cur.lastrowid, doc_id, meta["title"], ch["text"], norm_ch)
        result = {"status": status, "id": doc_id, "path": str(path), "engine": engine, "error": err}

    # Outside the DB lock: optional auto-filing of freshly indexed documents.
    if status == "indexed" and os.environ.get("DOCWISE_AUTO_FILE", "0") == "1":
        try:
            filed = auto_file_document(doc_id)
            if filed:
                result["auto_filed"] = filed
        except Exception:
            pass
    return result


def row_to_doc(row: sqlite3.Row, include_text: bool = False) -> dict:
    d = dict(row)
    d["tags"] = safe_json_load(d.get("tags"), [])
    d["fields"] = safe_json_load(d.get("fields"), {})
    if not include_text:
        d.pop("text", None)
        d.pop("normalized_text", None)
    return d


def scan_folder(path: Path, recursive: bool = True, force: bool = False) -> dict:
    path = path.expanduser().resolve()
    if not path.exists() or not path.is_dir():
        raise HTTPException(status_code=400, detail="Folder does not exist")
    pattern = "**/*" if recursive else "*"
    files = [p for p in path.glob(pattern) if p.is_file() and p.suffix.lower() in SUPPORTED]
    result = {"folder": str(path), "found": len(files), "indexed": 0, "unchanged": 0, "skipped": 0, "duplicates": 0, "errors": 0, "items": []}
    SCAN_STATE.update({"running": True, "message": f"Scanning {path}", "indexed": 0, "errors": 0})
    try:
        for p in files:
            try:
                item = index_file(p, "folder", force=force)
                result["items"].append(item)
                if item["status"] in ("indexed", "needs_ocr", "error"):
                    result["indexed"] += 1
                elif item["status"] == "unchanged":
                    result["unchanged"] += 1
                elif item["status"] == "duplicate":
                    result["duplicates"] += 1
                else:
                    result["skipped"] += 1
                if item["status"] == "error":
                    result["errors"] += 1
            except Exception as exc:
                result["errors"] += 1
                result["items"].append({"status": "error", "path": str(p), "error": str(exc)})
            SCAN_STATE.update({"indexed": result["indexed"], "errors": result["errors"]})
    finally:
        SCAN_STATE.update({"running": False, "last": now_iso(), "message": "Scan complete"})
    with DB_LOCK, db() as conn:
        conn.execute("UPDATE folders SET last_scan=? WHERE path=?", (now_iso(), str(path)))
    return result


def watcher_loop():
    while True:
        time.sleep(int(os.environ.get("DOCWISE_WATCH_INTERVAL", "45")))
        try:
            with DB_LOCK, db() as conn:
                folders = conn.execute("SELECT path, recursive FROM folders WHERE watch=1").fetchall()
            if SCAN_STATE["running"]:
                continue
            for f in folders:
                try:
                    scan_folder(Path(f["path"]), bool(f["recursive"]), force=False)
                except Exception as exc:
                    SCAN_STATE.update({"message": f"Watch scan error: {exc}", "errors": SCAN_STATE.get("errors", 0) + 1})
        except Exception:
            pass


threading.Thread(target=watcher_loop, daemon=True).start()
threading.Thread(target=embedding_maintenance_loop, daemon=True).start()


class FolderRequest(BaseModel):
    path: str
    recursive: bool = True
    watch: bool = True
    scan_now: bool = True


class SearchRequest(BaseModel):
    q: str = ""
    doc_type: str = "all"
    language: str = "all"
    limit: int = 50


class AskRequest(BaseModel):
    question: str
    limit: int = 6
    use_ai: bool = True


class OrganizeRequest(BaseModel):
    document_id: int
    mode: str = "copy"  # copy or move


class BulkDeleteRequest(BaseModel):
    ids: list[int] = []
    all: bool = False
    delete_files: bool = False
    app_files_only: bool = True


class TextUpdateRequest(BaseModel):
    text: str


class EvalCaseRequest(BaseModel):
    question: str
    expected: str = ""
    must_cite: str = ""


@app.get("/api/status")
def status():
    with DB_LOCK, db() as conn:
        total = conn.execute("SELECT COUNT(*) c FROM documents").fetchone()["c"]
        needs = conn.execute("SELECT COUNT(*) c FROM documents WHERE status!='indexed'").fetchone()["c"]
        folders = conn.execute("SELECT COUNT(*) c FROM folders").fetchone()["c"]
        chunks = conn.execute("SELECT COUNT(*) c FROM chunks").fetchone()["c"]
        embeddings = conn.execute("SELECT COUNT(*) c FROM chunk_embeddings WHERE dims>0").fetchone()["c"]
        fts_rows = conn.execute("SELECT COUNT(*) c FROM chunk_fts").fetchone()["c"]
        quality_rows = conn.execute("SELECT ocr_quality, COUNT(*) c FROM documents GROUP BY ocr_quality").fetchall()
        ocr_quality_counts = {r["ocr_quality"] or "unknown": r["c"] for r in quality_rows}
    auto_base = os.environ.get("DOCWISE_AUTO_FILE_BASE", "") or str(ARCHIVE)
    return {"ok": True, "documents": total, "needs_review": needs, "folders": folders, "chunks": chunks, "embeddings": embeddings, "fts_rows": fts_rows, "ocr_quality_counts": ocr_quality_counts, "scan": SCAN_STATE, "ocr_progress": dict(OCR_STATE), "auto_file": {"enabled": os.environ.get("DOCWISE_AUTO_FILE", "0") == "1", "base": auto_base}, "ocr": available_ocr()}


@app.post("/api/upload")
async def upload(files: list[UploadFile] = File(...)):
    results = []
    for f in files:
        ext = Path(f.filename or "file").suffix.lower()
        if ext not in SUPPORTED:
            results.append({"filename": f.filename, "status": "skipped", "reason": "unsupported"})
            continue
        stamp = datetime.now().strftime("%Y%m%d_%H%M%S_%f")
        safe = re.sub(r"[^A-Za-z0-9_.\-\u0600-\u06FF]+", "_", f.filename or f"upload{ext}")
        dest = UPLOADS / f"{stamp}_{safe}"
        with dest.open("wb") as out:
            shutil.copyfileobj(f.file, out)
        try:
            results.append(index_file(dest, "upload", force=True))
        except Exception as exc:
            results.append({"filename": f.filename, "status": "error", "error": str(exc)})
    return {"results": results}


@app.post("/api/pick-folder")
def pick_folder():
    """Open the native folder chooser on the machine running the app (the app
    is local-first, so that is the same machine as the browser)."""
    script = (
        "import tkinter as tk\n"
        "from tkinter import filedialog\n"
        "root = tk.Tk()\n"
        "root.withdraw()\n"
        "root.attributes('-topmost', True)\n"
        "root.update()\n"
        "print(filedialog.askdirectory(title='Choose a folder for DocWise to index') or '')\n"
    )
    try:
        proc = subprocess.run(
            [sys.executable, "-c", script], capture_output=True, text=True, timeout=300
        )
        chosen = (proc.stdout or "").strip()
        if proc.returncode != 0:
            raise RuntimeError((proc.stderr or "picker failed").strip()[:200])
        if not chosen:
            return {"path": None, "cancelled": True}
        return {"path": str(Path(chosen)), "cancelled": False}
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"Folder picker unavailable ({exc}). Type the path manually.")


@app.post("/api/folders")
def add_folder(req: FolderRequest):
    folder = Path(req.path).expanduser().resolve()
    if not folder.exists() or not folder.is_dir():
        raise HTTPException(status_code=400, detail="Folder does not exist")
    with DB_LOCK, db() as conn:
        conn.execute(
            "INSERT INTO folders(path, recursive, watch, created_at) VALUES(?,?,?,?) ON CONFLICT(path) DO UPDATE SET recursive=excluded.recursive, watch=excluded.watch",
            (str(folder), int(req.recursive), int(req.watch), now_iso()),
        )
    result = None
    if req.scan_now:
        result = scan_folder(folder, req.recursive, force=False)
    return {"folder": str(folder), "scan": result}


@app.get("/api/folders")
def list_folders():
    with DB_LOCK, db() as conn:
        rows = conn.execute("SELECT * FROM folders ORDER BY created_at DESC").fetchall()
    return {"folders": [dict(r) for r in rows]}


@app.delete("/api/folders/{folder_id}")
def delete_folder(folder_id: int):
    with DB_LOCK, db() as conn:
        conn.execute("DELETE FROM folders WHERE id=?", (folder_id,))
    return {"ok": True}


@app.post("/api/scan")
def manual_scan(req: FolderRequest):
    return scan_folder(Path(req.path), recursive=req.recursive, force=True)


@app.get("/api/documents")
def documents(limit: int = 80, offset: int = 0):
    with DB_LOCK, db() as conn:
        rows = conn.execute("SELECT * FROM documents ORDER BY updated_at DESC LIMIT ? OFFSET ?", (limit, offset)).fetchall()
    return {"documents": [row_to_doc(r) for r in rows]}


@app.get("/api/documents/{doc_id}")
def document(doc_id: int):
    with DB_LOCK, db() as conn:
        row = conn.execute("SELECT * FROM documents WHERE id=?", (doc_id,)).fetchone()
        chunks = conn.execute("SELECT page, chunk_index, text FROM chunks WHERE document_id=? ORDER BY page, chunk_index", (doc_id,)).fetchall()
    if not row:
        raise HTTPException(status_code=404, detail="Document not found")
    d = row_to_doc(row, include_text=True)
    d["chunks"] = [dict(c) for c in chunks]
    d["suggestion"] = filing_suggestion(d)
    return d


class OcrEngineRequest(BaseModel):
    engine: str = "auto"


def persist_env_line(key: str, value: str) -> None:
    """Best-effort persistence of one setting into .env (runtime env is
    already set by the caller)."""
    try:
        env_path = ROOT / ".env"
        lines = env_path.read_text(encoding="utf-8").splitlines() if env_path.exists() else []
        lines = [ln for ln in lines if not re.match(rf"\s*#?\s*{re.escape(key)}\s*=", ln)]
        lines.append(f"{key}={value}")
        env_path.write_text("\n".join(lines) + "\n", encoding="utf-8")
    except Exception:
        pass


@app.post("/api/settings/ocr-engine")
def set_ocr_engine(req: OcrEngineRequest):
    """Choose which OCR engine acts: auto (confidence-based chain) or one
    forced engine. Applies immediately and persists to .env."""
    engine = (req.engine or "auto").lower()
    allowed = {"auto", "tesseract", "ollama", "azure", "openai"}
    if engine not in allowed:
        raise HTTPException(status_code=400, detail=f"engine must be one of {sorted(allowed)}")
    os.environ["DOCWISE_OCR"] = engine
    persist_env_line("DOCWISE_OCR", engine)
    return {"ok": True, "engine": engine}


class AutoFileRequest(BaseModel):
    enabled: bool = False


@app.post("/api/settings/auto-file")
def set_auto_file(req: AutoFileRequest):
    """One-step mode: newly indexed documents are auto-copied into the
    archive filing tree, no confirmation. Copy only - originals stay."""
    value = "1" if req.enabled else "0"
    os.environ["DOCWISE_AUTO_FILE"] = value
    persist_env_line("DOCWISE_AUTO_FILE", value)
    return {"ok": True, "enabled": req.enabled}


@app.get("/api/export/xlsx")
def export_xlsx():
    """All documents with extracted fields as a spreadsheet."""
    if not Workbook:
        raise HTTPException(status_code=400, detail="openpyxl is not installed")
    with DB_LOCK, db() as conn:
        rows = conn.execute("SELECT * FROM documents ORDER BY updated_at DESC").fetchall()
    wb = Workbook()
    ws = wb.active
    ws.title = "Documents"
    headers = ["ID", "Title", "Type", "Language", "Date", "Company", "Amount", "Total", "Currency",
               "Invoice No", "VAT No", "IBAN", "Due/Expiry", "Status", "OCR Quality", "File", "Summary"]
    ws.append(headers)
    for r in rows:
        d = row_to_doc(r)
        f = d.get("fields") or {}
        ws.append([
            d.get("id"), d.get("title"), d.get("doc_type"), d.get("language"),
            d.get("date_guess"), d.get("company"), d.get("amount"),
            f.get("total"), f.get("currency"), f.get("invoice_number"),
            f.get("vat_number"), f.get("iban"), f.get("due_or_expiry_date"),
            d.get("status"), d.get("ocr_quality"), d.get("path"),
            (d.get("summary") or "")[:500],
        ])
    for col, width in zip("ABCDEFGHIJKLMNOPQ", (6, 34, 12, 16, 14, 24, 12, 12, 10, 16, 16, 26, 16, 12, 12, 50, 60)):
        ws.column_dimensions[col].width = width
    derived = DATA / "derived"
    derived.mkdir(parents=True, exist_ok=True)
    out = derived / "docwise_export.xlsx"
    wb.save(str(out))
    return FileResponse(str(out), filename="docwise_documents.xlsx",
                        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")


@app.get("/api/documents/{doc_id}/page/{page_num}/image")
def page_image(doc_id: int, page_num: int):
    """Rendered page image for the viewer (PDF pages render on demand)."""
    with DB_LOCK, db() as conn:
        row = conn.execute("SELECT path, file_ext FROM documents WHERE id=?", (doc_id,)).fetchone()
    if not row:
        raise HTTPException(status_code=404, detail="Document not found")
    path = Path(row["path"])
    if not path.exists():
        raise HTTPException(status_code=404, detail="File not found")
    ext = (row["file_ext"] or "").lower()
    if ext in IMAGE_EXTS:
        return FileResponse(str(path))
    if ext == ".pdf" and fitz:
        try:
            doc = fitz.open(path)
            if page_num < 1 or page_num > doc.page_count:
                raise HTTPException(status_code=404, detail="Page out of range")
            pix = doc[page_num - 1].get_pixmap(matrix=fitz.Matrix(2.0, 2.0), alpha=False)
            return Response(content=pix.tobytes("png"), media_type="image/png")
        except HTTPException:
            raise
        except Exception as exc:
            raise HTTPException(status_code=500, detail=str(exc))
    raise HTTPException(status_code=400, detail="No page image for this file type")


@app.get("/api/documents/{doc_id}/page/{page_num}/words")
def page_word_boxes(doc_id: int, page_num: int, q: str = ""):
    """OCR word boxes (page-relative fractions) with indexes matching query q."""
    with DB_LOCK, db() as conn:
        row = conn.execute("SELECT words FROM page_words WHERE document_id=? AND page=?", (doc_id, page_num)).fetchone()
        pages = conn.execute(
            "SELECT DISTINCT page FROM page_words WHERE document_id=? ORDER BY page", (doc_id,)
        ).fetchall()
    words = safe_json_load(row["words"], []) if row else []
    matched = []
    terms = [t for t in normalize_arabic(q).split() if len(t) > 1] if q.strip() else []
    if terms:
        for i, w in enumerate(words):
            nw = normalize_arabic(w.get("t", ""))
            if nw and any(t in nw or (len(nw) > 2 and nw in t) for t in terms):
                matched.append(i)
    return {"page": page_num, "words": words, "matched": matched, "pages_with_words": [p["page"] for p in pages]}


@app.get("/api/documents/{doc_id}/searchable-pdf")
def searchable_pdf(doc_id: int):
    """Export scans as a real searchable PDF: original pages with an invisible
    Tesseract text layer on pages that had no text."""
    with DB_LOCK, db() as conn:
        row = conn.execute("SELECT * FROM documents WHERE id=?", (doc_id,)).fetchone()
    if not row:
        raise HTTPException(status_code=404, detail="Document not found")
    path = Path(row["path"])
    if not path.exists():
        raise HTTPException(status_code=404, detail="File not found")
    ext = (row["file_ext"] or "").lower()
    if not (fitz and pytesseract and Image and tesseract_cmd()):
        raise HTTPException(status_code=400, detail="Searchable PDF export needs PyMuPDF and Tesseract")
    derived = DATA / "derived"
    derived.mkdir(parents=True, exist_ok=True)
    out_path = derived / f"{doc_id}_searchable.pdf"
    download_name = f"{path.stem}_searchable.pdf"
    if out_path.exists() and out_path.stat().st_mtime >= (row["mtime"] or 0):
        return FileResponse(str(out_path), filename=download_name, media_type="application/pdf")
    pytesseract.pytesseract.tesseract_cmd = tesseract_cmd()
    td = tessdata_dir()
    if td:
        os.environ["TESSDATA_PREFIX"] = td
    langs = set(available_languages())
    lang = "ara+eng" if {"ara", "eng"}.issubset(langs) else ("ara" if "ara" in langs else "eng")
    max_pages = int(os.environ.get("DOCWISE_MAX_OCR_PDF_PAGES", "20"))
    out = fitz.open()
    try:
        if ext in IMAGE_EXTS:
            pdf_bytes = pytesseract.image_to_pdf_or_hocr(Image.open(path), extension="pdf", lang=lang)
            with fitz.open("pdf", pdf_bytes) as page_doc:
                out.insert_pdf(page_doc)
        elif ext == ".pdf":
            with fitz.open(path) as src:
                for i, page in enumerate(src, start=1):
                    text = page.get_text("text") or ""
                    if not should_ocr_pdf_page(text) or i > max_pages:
                        out.insert_pdf(src, from_page=i - 1, to_page=i - 1)
                        continue
                    pix = page.get_pixmap(matrix=fitz.Matrix(3.5, 3.5), alpha=False)
                    img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
                    pdf_bytes = pytesseract.image_to_pdf_or_hocr(img, extension="pdf", lang=lang)
                    with fitz.open("pdf", pdf_bytes) as page_doc:
                        out.insert_pdf(page_doc)
        else:
            raise HTTPException(status_code=400, detail="Searchable PDF export works for PDFs and images")
        out.save(str(out_path))
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"Searchable PDF failed: {exc}")
    finally:
        out.close()
    return FileResponse(str(out_path), filename=download_name, media_type="application/pdf")


def is_app_managed_file(path: Path) -> bool:
    try:
        resolved = path.resolve()
        return resolved.is_relative_to(UPLOADS.resolve()) or resolved.is_relative_to(ARCHIVE.resolve())
    except Exception:
        return False


def delete_docs_by_ids(ids: list[int], delete_files: bool = False, app_files_only: bool = True) -> dict:
    ids = sorted(set(int(i) for i in ids if int(i) > 0))
    if not ids:
        return {"ok": True, "deleted_count": 0, "file_deleted_count": 0, "items": []}
    placeholders = ",".join("?" for _ in ids)
    items = []
    file_deleted_count = 0
    with DB_LOCK, db() as conn:
        rows = conn.execute(f"SELECT id, path FROM documents WHERE id IN ({placeholders})", ids).fetchall()
        for row in rows:
            path = Path(row["path"])
            item = {"id": row["id"], "path": str(path), "index_deleted": True, "file_deleted": False, "file_skipped": False, "file_error": None}
            if delete_files:
                if app_files_only and not is_app_managed_file(path):
                    item["file_skipped"] = True
                    item["file_error"] = "Skipped external file because app_files_only=true"
                else:
                    try:
                        if path.exists() and path.is_file():
                            path.unlink()
                            item["file_deleted"] = True
                            file_deleted_count += 1
                    except Exception as exc:
                        item["file_error"] = str(exc)
            items.append(item)
        for doc_id in ids:
            delete_chunk_indexes_for_doc(conn, doc_id)
        conn.execute(f"DELETE FROM chunks WHERE document_id IN ({placeholders})", ids)
        conn.execute(f"DELETE FROM page_words WHERE document_id IN ({placeholders})", ids)
        conn.execute(f"DELETE FROM documents WHERE id IN ({placeholders})", ids)
    return {"ok": True, "deleted_count": len(items), "file_deleted_count": file_deleted_count, "items": items}


@app.put("/api/documents/{doc_id}/text")
def update_document_text(doc_id: int, req: TextUpdateRequest):
    """Save manually corrected OCR text and rebuild search/RAG chunks."""
    text = req.text or ""
    with DB_LOCK, db() as conn:
        row = conn.execute("SELECT * FROM documents WHERE id=?", (doc_id,)).fetchone()
        if not row:
            raise HTTPException(status_code=404, detail="Document not found")
        path = Path(row["path"])
        meta = guess_metadata(path, text)
        fields = extract_structured_fields(text, meta)
        quality, quality_score = ocr_quality(text)
        normalized = normalize_arabic(text + " " + meta["title"] + " " + " ".join(meta["tags"]) + " " + json.dumps(fields, ensure_ascii=False))
        status = "indexed" if text.strip() else "needs_ocr"
        conn.execute(
            """
            UPDATE documents SET title=?, doc_type=?, language=?, date_guess=?, company=?, amount=?, tags=?, summary=?, fields=?,
                text=?, normalized_text=?, status=?, ocr_engine=?, ocr_quality=?, ocr_score=?, error=?, updated_at=? WHERE id=?
            """,
            (
                meta["title"], meta["doc_type"], meta["language"], meta["date_guess"], meta["company"], meta["amount"],
                json.dumps(meta["tags"], ensure_ascii=False), meta["summary"], json.dumps(fields, ensure_ascii=False), text, normalized, status,
                "manual-correction", quality, quality_score, None, now_iso(), doc_id,
            ),
        )
        delete_chunk_indexes_for_doc(conn, doc_id)
        conn.execute("DELETE FROM chunks WHERE document_id=?", (doc_id,))
        for ch in chunk_pages([{"page": 1, "text": text}]):
            norm_ch = normalize_arabic(ch["text"])
            cur = conn.execute(
                "INSERT INTO chunks(document_id, page, chunk_index, text, normalized_text, token_count) VALUES(?,?,?,?,?,?)",
                (doc_id, ch["page"], ch["chunk_index"], ch["text"], norm_ch, ch.get("token_count") or approx_token_count(ch["text"])),
            )
            add_chunk_indexes(conn, cur.lastrowid, doc_id, meta["title"], ch["text"], norm_ch)
    return {"ok": True, "id": doc_id, "status": status}


@app.post("/api/documents/{doc_id}/vision-ocr")
def force_vision_ocr(doc_id: int):
    """Force OpenAI Vision OCR for one image/PDF and rebuild metadata/RAG indexes."""
    if not (os.environ.get("OPENAI_API_KEY") and OpenAI):
        raise HTTPException(status_code=400, detail="OPENAI_API_KEY is not configured")
    with DB_LOCK, db() as conn:
        row = conn.execute("SELECT path FROM documents WHERE id=?", (doc_id,)).fetchone()
    if not row:
        raise HTTPException(status_code=404, detail="Document not found")
    path = Path(row["path"])
    if not path.exists():
        raise HTTPException(status_code=404, detail="Original file not found")
    ext = path.suffix.lower()
    pages = []
    engine = "openai-vision"
    errors = []
    if ext in IMAGE_EXTS:
        text, used = ocr_image_openai(path)
        engine = used
        pages = [{"page": 1, "text": text}]
        if not text.strip():
            errors.append(used)
    elif ext == ".pdf":
        if not fitz:
            raise HTTPException(status_code=400, detail="PyMuPDF is not installed")
        doc = fitz.open(path)
        max_pages = int(os.environ.get("DOCWISE_MAX_OCR_PDF_PAGES", "20"))
        for i, page in enumerate(doc, start=1):
            if i > max_pages:
                break
            pix = page.get_pixmap(matrix=fitz.Matrix(2.5, 2.5), alpha=False)
            with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                tmp_path = Path(tmp.name)
            pix.save(str(tmp_path))
            text, used = ocr_image_openai(tmp_path)
            tmp_path.unlink(missing_ok=True)
            pages.append({"page": i, "text": text})
            engine = f"pdf-render+{used}"
            if not text.strip():
                errors.append(f"page {i}: {used}")
    else:
        raise HTTPException(status_code=400, detail="Vision OCR supports images and PDFs only")

    full_text = "\n\n".join(p.get("text") or "" for p in pages).strip()
    meta = guess_metadata(path, full_text)
    fields = extract_structured_fields(full_text, meta)
    quality, quality_score = ocr_quality(full_text)
    normalized = normalize_arabic(full_text + " " + meta["title"] + " " + " ".join(meta["tags"]) + " " + json.dumps(fields, ensure_ascii=False))
    status = "indexed" if full_text else "needs_ocr"
    err = "; ".join(errors) if errors else None
    with DB_LOCK, db() as conn:
        conn.execute(
            """
            UPDATE documents SET title=?, doc_type=?, language=?, date_guess=?, company=?, amount=?, tags=?, summary=?, fields=?,
                text=?, normalized_text=?, status=?, ocr_engine=?, ocr_quality=?, ocr_score=?, error=?, updated_at=? WHERE id=?
            """,
            (meta["title"], meta["doc_type"], meta["language"], meta["date_guess"], meta["company"], meta["amount"],
             json.dumps(meta["tags"], ensure_ascii=False), meta["summary"], json.dumps(fields, ensure_ascii=False),
             full_text, normalized, status, engine, quality, quality_score, err, now_iso(), doc_id),
        )
        delete_chunk_indexes_for_doc(conn, doc_id)
        conn.execute("DELETE FROM chunks WHERE document_id=?", (doc_id,))
        for ch in chunk_pages(pages):
            norm_ch = normalize_arabic(ch["text"])
            cur = conn.execute(
                "INSERT INTO chunks(document_id, page, chunk_index, text, normalized_text, token_count) VALUES(?,?,?,?,?,?)",
                (doc_id, ch["page"], ch["chunk_index"], ch["text"], norm_ch, ch.get("token_count") or approx_token_count(ch["text"])),
            )
            add_chunk_indexes(conn, cur.lastrowid, doc_id, meta["title"], ch["text"], norm_ch)
    return {"ok": True, "id": doc_id, "status": status, "engine": engine, "ocr_quality": quality, "ocr_score": quality_score, "error": err}


@app.delete("/api/documents/{doc_id}")
def delete_document(doc_id: int):
    """Remove only the DocWise index record. The original file stays untouched."""
    result = delete_docs_by_ids([doc_id], delete_files=False)
    result["deleted"] = "index_only"
    return result


@app.delete("/api/documents/{doc_id}/file")
def delete_document_and_file(doc_id: int):
    """Delete the original file from disk, then remove the DocWise index record."""
    result = delete_docs_by_ids([doc_id], delete_files=True, app_files_only=False)
    result["deleted"] = "file_and_index"
    return result


@app.post("/api/documents/bulk-delete")
def bulk_delete_documents(req: BulkDeleteRequest):
    with DB_LOCK, db() as conn:
        if req.all:
            ids = [r["id"] for r in conn.execute("SELECT id FROM documents").fetchall()]
        else:
            ids = req.ids
    return delete_docs_by_ids(ids, delete_files=req.delete_files, app_files_only=req.app_files_only)


@app.post("/api/reindex/{doc_id}")
def reindex(doc_id: int):
    with DB_LOCK, db() as conn:
        row = conn.execute("SELECT path FROM documents WHERE id=?", (doc_id,)).fetchone()
    if not row:
        raise HTTPException(status_code=404, detail="Document not found")
    return index_file(Path(row["path"]), force=True)


@app.post("/api/reclassify")
def reclassify_all():
    """Re-run classification and field extraction on the stored text of every
    document - no re-OCR, so it is fast and free. Use after rule updates."""
    updated = 0
    with DB_LOCK, db() as conn:
        rows = conn.execute("SELECT id, path, text FROM documents").fetchall()
        for r in rows:
            text = r["text"] or ""
            meta = guess_metadata(Path(r["path"]), text)
            fields = extract_structured_fields(text, meta)
            normalized = normalize_arabic(text + " " + meta["title"] + " " + " ".join(meta["tags"]) + " " + json.dumps(fields, ensure_ascii=False))
            conn.execute(
                """UPDATE documents SET title=?, doc_type=?, language=?, date_guess=?, company=?, amount=?,
                   tags=?, summary=?, fields=?, normalized_text=?, updated_at=? WHERE id=?""",
                (meta["title"], meta["doc_type"], meta["language"], meta["date_guess"], meta["company"], meta["amount"],
                 json.dumps(meta["tags"], ensure_ascii=False), meta["summary"], json.dumps(fields, ensure_ascii=False),
                 normalized, now_iso(), r["id"]),
            )
            updated += 1
    return {"ok": True, "reclassified": updated}


@app.post("/api/rebuild-rag")
def rebuild_rag_indexes():
    """Rebuild FTS/vector indexes for already extracted documents without re-OCRing files."""
    rebuilt = 0
    with DB_LOCK, db() as conn:
        docs = conn.execute("SELECT id, title FROM documents WHERE status='indexed'").fetchall()
        for doc in docs:
            delete_chunk_indexes_for_doc(conn, doc["id"])
            chunks = conn.execute("SELECT * FROM chunks WHERE document_id=? ORDER BY page, chunk_index", (doc["id"],)).fetchall()
            for ch in chunks:
                add_chunk_indexes(conn, ch["id"], doc["id"], doc["title"], ch["text"], ch["normalized_text"])
                rebuilt += 1
    return {"ok": True, "rebuilt_chunks": rebuilt}


@app.post("/api/search")
def search(req: SearchRequest):
    tokens = query_terms(req.q)
    if req.q.strip():
        retrieved = hybrid_retrieve(req.q, limit=max(30, req.limit * 3), doc_type=None if req.doc_type == "all" else req.doc_type)
        doc_map = {}
        for ch in retrieved:
            doc_id = ch["document_id"]
            if doc_id not in doc_map or ch["score"] > doc_map[doc_id]["score"]:
                doc_map[doc_id] = ch
        ids = list(doc_map.keys())[: max(1, min(req.limit, 100))]
        if not ids:
            return {"results": []}
        placeholders = ",".join("?" for _ in ids)
        with DB_LOCK, db() as conn:
            rows = conn.execute(f"SELECT * FROM documents WHERE id IN ({placeholders})", ids).fetchall()
        by_id = {r["id"]: row_to_doc(r, include_text=True) for r in rows}
        results = []
        for doc_id in ids:
            d = by_id.get(doc_id)
            if not d:
                continue
            if req.language != "all" and req.language not in (d.get("language") or ""):
                continue
            best = doc_map[doc_id]
            d["score"] = round(best["score"], 4)
            d["snippet"] = make_snippet(best.get("text") or d.get("summary") or "", tokens)
            d["retrieval"] = best.get("retrieval", "hybrid")
            d.pop("text", None)
            d.pop("normalized_text", None)
            results.append(d)
        return {"results": results}

    with DB_LOCK, db() as conn:
        sql = "SELECT * FROM documents WHERE 1=1"
        params = []
        if req.doc_type != "all":
            sql += " AND doc_type=?"
            params.append(req.doc_type)
        if req.language != "all":
            sql += " AND language LIKE ?"
            params.append(f"%{req.language}%")
        rows = conn.execute(sql + " ORDER BY updated_at DESC LIMIT ?", params + [max(1, min(req.limit, 100))]).fetchall()
    return {"results": [row_to_doc(r) for r in rows]}


def make_snippet(text: str, tokens: list[str]) -> str:
    clean = re.sub(r"\s+", " ", text or "").strip()
    if not clean:
        return ""
    norm = normalize_arabic(clean)
    idx = 0
    for t in tokens:
        pos = norm.find(t)
        if pos >= 0:
            idx = max(0, min(len(clean), pos) - 100)
            break
    return clean[idx:idx + 360] + ("..." if len(clean) > idx + 360 else "")


def query_terms(text: str) -> list[str]:
    base = [t for t in normalize_arabic(text).split() if len(t) > 1]
    synonyms = {
        "فاتوره": ["invoice", "bill", "total", "amount", "aed", "درهم", "اجمالي", "مبلغ"],
        "الفاتوره": ["فاتوره", "invoice", "bill", "total", "amount", "aed", "درهم", "اجمالي", "مبلغ"],
        "قيمه": ["amount", "total", "مبلغ", "اجمالي", "aed", "درهم"],
        "المبلغ": ["amount", "total", "مبلغ", "اجمالي", "aed", "درهم"],
        "كم": ["amount", "total", "مبلغ", "اجمالي"],
        "عقد": ["contract", "agreement", "lease", "rent", "ايجار"],
        "هويه": ["identity", "passport", "emirates", "id", "expiry", "انتهاء"],
        "تاريخ": ["date", "expiry", "due", "تاريخ"],
    }
    out = []
    for t in base:
        out.append(t)
        for prefix in ("ال", "و", "ب", "بال", "لل"):
            if t.startswith(prefix) and len(t) > len(prefix) + 1:
                out.append(t[len(prefix):])
        out.extend(synonyms.get(t, []))
    seen = set()
    uniq = []
    for t in out:
        nt = normalize_arabic(t)
        if nt and nt not in seen:
            seen.add(nt)
            uniq.append(nt)
    return uniq


def chunk_score(question: str, chunk_norm: str) -> int:
    tokens = query_terms(question)
    if not tokens:
        return 0
    words = set(chunk_norm.split())
    score = 0
    for t in tokens:
        if t in words:
            score += 4
        if t in chunk_norm:
            score += chunk_norm.count(t)
        if len(t) > 3 and any(w.startswith(t) or t.startswith(w) for w in words if len(w) > 3):
            score += 1
    return score


def rerank_candidate(candidate: dict, question: str, plan: dict) -> float:
    score = float(candidate.get("score", 0.0))
    text = candidate.get("text") or ""
    norm = candidate.get("normalized_text") or normalize_arabic(text)
    score += min(3.0, chunk_score(question, norm) * 0.08)
    if plan.get("doc_type") and candidate.get("doc_type") == plan["doc_type"]:
        score += 1.2
    if plan.get("field") == "amount" and re.search(r"(?:AED|د\.إ|درهم|total|amount|اجمالي|مبلغ|\d+[,.]?\d*)", text, re.I):
        score += 0.9
    if plan.get("field") == "date" and re.search(r"20\d{2}|\d{1,2}[-/.]\d{1,2}|expiry|due|انتهاء|تاريخ", text, re.I):
        score += 0.7
    if plan.get("sort") == "latest" and candidate.get("updated_at"):
        score += 0.15
    return score


def gpt_rerank(question: str, candidates: list[dict], limit: int) -> list[dict]:
    if not (os.environ.get("OPENAI_API_KEY") and OpenAI) or os.environ.get("DOCWISE_GPT_RERANK", "1") != "1":
        return candidates[:limit]
    try:
        client = OpenAI()
        packed = []
        for i, c in enumerate(candidates[:20], start=1):
            packed.append({"n": i, "title": c.get("title"), "type": c.get("doc_type"), "page": c.get("page"), "text": (c.get("text") or "")[:900]})
        resp = client.chat.completions.create(
            model=os.environ.get("DOCWISE_CHAT_MODEL", "gpt-4o-mini"),
            messages=[
                {"role": "system", "content": "Rerank document chunks for answering the user question. Return JSON only: {\"order\":[candidate_numbers_best_first]}. Prefer exact supporting evidence."},
                {"role": "user", "content": json.dumps({"question": question, "candidates": packed}, ensure_ascii=False)},
            ],
            temperature=0,
        )
        content = (resp.choices[0].message.content or "{}").strip()
        content = re.sub(r"^```json|```$", "", content, flags=re.I).strip()
        order = json.loads(content).get("order", [])
        by_num = {i: c for i, c in enumerate(candidates[:20], start=1)}
        reranked = [by_num[n] for n in order if n in by_num]
        seen = {id(c) for c in reranked}
        reranked.extend([c for c in candidates if id(c) not in seen])
        for c in reranked[:limit]:
            c["retrieval"] = (c.get("retrieval") or "hybrid") + "+gpt-rerank"
        return reranked[:limit]
    except Exception:
        return candidates[:limit]


def verify_answer(question: str, answer: str, sources: list[dict]) -> dict:
    if not (os.environ.get("OPENAI_API_KEY") and OpenAI) or os.environ.get("DOCWISE_VERIFY", "1") != "1":
        return {"enabled": False}
    try:
        client = OpenAI()
        resp = client.chat.completions.create(
            model=os.environ.get("DOCWISE_CHAT_MODEL", "gpt-4o-mini"),
            messages=[
                {"role": "system", "content": "Verify if the answer is fully supported by the sources. Return JSON only with keys supported:boolean, confidence:0-1, reason:string."},
                {"role": "user", "content": json.dumps({"question": question, "answer": answer, "sources": sources}, ensure_ascii=False)},
            ],
            temperature=0,
        )
        content = (resp.choices[0].message.content or "{}").strip()
        content = re.sub(r"^```json|```$", "", content, flags=re.I).strip()
        data = json.loads(content)
        data["enabled"] = True
        return data
    except Exception as exc:
        return {"enabled": True, "error": str(exc)[:160]}


def hybrid_retrieve(question: str, limit: int = 8, doc_type: Optional[str] = None) -> list[dict]:
    """Hybrid RAG retrieval: SQLite FTS5/BM25 + vector embeddings + Arabic-aware lexical reranking."""
    plan = query_plan(question)
    if doc_type:
        plan["doc_type"] = doc_type
    terms = plan.get("terms") or query_terms(question)
    fts_query = fts_query_from_terms(terms)
    merged: dict[int, dict] = {}

    with DB_LOCK, db() as conn:
        params_base = []
        doc_filter = ""
        if plan.get("doc_type"):
            doc_filter = " AND documents.doc_type=?"
            params_base.append(plan["doc_type"])

        if fts_query:
            try:
                rows = conn.execute(
                    f"""
                    SELECT chunks.*, documents.title, documents.path, documents.doc_type, documents.updated_at,
                           documents.date_guess, documents.amount, documents.company, bm25(chunk_fts) AS bm25_rank
                    FROM chunk_fts
                    JOIN chunks ON chunks.id=chunk_fts.chunk_id
                    JOIN documents ON documents.id=chunks.document_id
                    WHERE chunk_fts MATCH ? AND documents.status='indexed'{doc_filter}
                    ORDER BY bm25_rank LIMIT 60
                    """,
                    [fts_query] + params_base,
                ).fetchall()
                for r in rows:
                    d = dict(r)
                    # FTS5 bm25 is usually negative; more negative is better.
                    d["score"] = 2.0 + max(0.0, -float(d.get("bm25_rank") or 0.0))
                    d["retrieval"] = "fts5-bm25"
                    merged[d["id"]] = d
            except Exception:
                pass

        # Vector retrieval. OpenAI embeddings if configured, local ONNX e5 when
        # ready, hash vectors as last resort. Only compare vectors built by the
        # same model as the query: cross-model cosine is meaningless noise.
        try:
            q_emb, q_model = embed_text(question, kind="query")
            rows = conn.execute(
                f"""
                SELECT chunks.*, documents.title, documents.path, documents.doc_type, documents.updated_at,
                       documents.date_guess, documents.amount, documents.company, chunk_embeddings.embedding,
                       chunk_embeddings.dims, chunk_embeddings.model
                FROM chunk_embeddings
                JOIN chunks ON chunks.id=chunk_embeddings.chunk_id
                JOIN documents ON documents.id=chunks.document_id
                WHERE chunk_embeddings.dims>0 AND chunk_embeddings.model=? AND documents.status='indexed'{doc_filter}
                """,
                [q_model] + params_base,
            ).fetchall()
            vector_hits = []
            for r in rows:
                try:
                    emb = json.loads(r["embedding"])
                    sim = cosine_similarity(q_emb, emb)
                    model = str(r["model"])
                    if model == ONNX_EMBED_MODEL:
                        # e5 cosine clusters around 0.70-0.92 even for unrelated
                        # text; rescale so real matches keep ranking spread.
                        if sim <= 0.78:
                            continue
                        score = 1.5 + (sim - 0.75) * 10.0
                    else:
                        min_sim = 0.18 if model.startswith("local-hash") else 0.12
                        if sim <= min_sim:
                            continue
                        score = 1.5 + sim * 3.0
                    d = dict(r)
                    d.pop("embedding", None)
                    d["score"] = score
                    d["retrieval"] = f"vector:{model}"
                    vector_hits.append(d)
                except Exception:
                    continue
            vector_hits.sort(key=lambda x: x["score"], reverse=True)
            for d in vector_hits[:60]:
                if d["id"] in merged:
                    merged[d["id"]]["score"] += d["score"]
                    merged[d["id"]]["retrieval"] += "+vector"
                else:
                    merged[d["id"]] = d
        except Exception:
            pass

        # Safety fallback: lexical scan if indexes are empty/stale.
        if not merged:
            rows = conn.execute(
                f"""
                SELECT chunks.*, documents.title, documents.path, documents.doc_type, documents.updated_at,
                       documents.date_guess, documents.amount, documents.company
                FROM chunks JOIN documents ON documents.id=chunks.document_id
                WHERE documents.status='indexed'{doc_filter}
                """,
                params_base,
            ).fetchall()
            for r in rows:
                s = chunk_score(question, r["normalized_text"])
                if s > 0:
                    d = dict(r)
                    d["score"] = s / 5.0
                    d["retrieval"] = "lexical-fallback"
                    merged[d["id"]] = d

    ranked = []
    for d in merged.values():
        d["score"] = rerank_candidate(d, question, plan)
        ranked.append(d)
    ranked.sort(key=lambda x: x["score"], reverse=True)
    return gpt_rerank(question, ranked, max(1, min(limit, 20)))


@app.get("/api/eval-cases")
def list_eval_cases():
    with DB_LOCK, db() as conn:
        rows = conn.execute("SELECT * FROM eval_cases ORDER BY id DESC").fetchall()
    return {"cases": [dict(r) for r in rows]}


@app.post("/api/eval-cases")
def add_eval_case(req: EvalCaseRequest):
    with DB_LOCK, db() as conn:
        cur = conn.execute(
            "INSERT INTO eval_cases(question, expected, must_cite, created_at) VALUES(?,?,?,?)",
            (req.question, req.expected, req.must_cite, now_iso()),
        )
    return {"ok": True, "id": cur.lastrowid}


@app.delete("/api/eval-cases/{case_id}")
def delete_eval_case(case_id: int):
    with DB_LOCK, db() as conn:
        conn.execute("DELETE FROM eval_cases WHERE id=?", (case_id,))
    return {"ok": True}


@app.post("/api/evaluate-rag")
def evaluate_rag():
    with DB_LOCK, db() as conn:
        cases = conn.execute("SELECT * FROM eval_cases ORDER BY id").fetchall()
    results = []
    for c in cases:
        answer_data = ask(AskRequest(question=c["question"], use_ai=bool(os.environ.get("OPENAI_API_KEY"))))
        answer = answer_data.get("answer", "")
        sources = answer_data.get("sources", [])
        expected_ok = not c["expected"] or normalize_arabic(c["expected"]) in normalize_arabic(answer)
        cite_ok = not c["must_cite"] or any(normalize_arabic(c["must_cite"]) in normalize_arabic(s.get("title", "") + " " + s.get("path", "")) for s in sources)
        results.append({
            "id": c["id"], "question": c["question"], "expected": c["expected"], "must_cite": c["must_cite"],
            "answer": answer, "sources": sources, "expected_ok": expected_ok, "cite_ok": cite_ok,
            "pass": expected_ok and cite_ok,
        })
    passed = sum(1 for r in results if r["pass"])
    return {"ok": True, "total": len(results), "passed": passed, "score": round(passed / len(results), 3) if results else None, "results": results}


@app.post("/api/ask")
def ask(req: AskRequest):
    top = hybrid_retrieve(req.question, limit=max(1, min(req.limit, 12)))
    if not top:
        return {"answer": "No matching indexed text found yet. Add folders/uploads and make sure OCR is working.", "sources": [], "mode": "hybrid-empty"}

    sources = [{"document_id": t["document_id"], "title": t["title"], "page": t["page"], "path": t["path"], "snippet": t["text"][:650], "score": round(t.get("score", 0), 4), "retrieval": t.get("retrieval", "hybrid")} for t in top]
    # Detect the question language in code: small local models otherwise answer
    # in the language of the sources, not the question.
    answer_lang = "أجب باللغة العربية فقط. Answer in Arabic ONLY." if has_arabic(req.question) else "Answer in English ONLY, even though the sources may be Arabic."
    if req.use_ai and not (os.environ.get("OPENAI_API_KEY") and OpenAI):
        # No OpenAI key: answer with a local Ollama model when one is running,
        # so Ask reasons over the sources instead of just listing excerpts.
        context = "\n\n".join([f"[Source {i+1}: {t['title']} page {t['page']}]\n{t['text'][:2500]}" for i, t in enumerate(top)])
        local = ollama_answer(req.question, context, answer_lang)
        if local:
            return {"answer": local, "sources": sources, "mode": f"ai-local:{ollama_chat_model()}"}
    if req.use_ai and os.environ.get("OPENAI_API_KEY") and OpenAI:
        try:
            client = OpenAI()
            context = "\n\n".join([f"[Source {i+1}: {t['title']} page {t['page']}]\n{t['text']}" for i, t in enumerate(top)])
            resp = client.chat.completions.create(
                model=os.environ.get("DOCWISE_CHAT_MODEL", "gpt-4o-mini"),
                messages=[
                    {"role": "system", "content": f"You are a strict document RAG assistant. Use only the provided sources. Cite sources like [Source 1]. If the sources do not clearly support the answer, say you could not find it in the indexed documents. Prefer exact amounts, dates, names, and quotes from the sources. {answer_lang}"},
                    {"role": "user", "content": f"Question: {req.question}\n\nSources:\n{context}"},
                ],
                temperature=0.2,
            )
            answer = resp.choices[0].message.content or ""
            verification = verify_answer(req.question, answer, sources)
            return {"answer": answer, "sources": sources, "mode": "ai", "verification": verification}
        except Exception as exc:
            return {"answer": fallback_answer(req.question, top) + f"\n\nAI error: {exc}", "sources": sources, "mode": "extractive"}
    return {"answer": fallback_answer(req.question, top), "sources": sources, "mode": "extractive"}


def fallback_answer(question: str, top: list[dict]) -> str:
    arabic = has_arabic(question)
    if arabic:
        intro = "أقرب النتائج من الأرشيف:"
        lines = [intro]
        for i, t in enumerate(top[:4], start=1):
            lines.append(f"[{i}] {t['title']} - صفحة {t['page']}: {t['text'][:360]}...")
        return "\n".join(lines)
    lines = ["Closest matches from your archive:"]
    for i, t in enumerate(top[:4], start=1):
        lines.append(f"[{i}] {t['title']} - page {t['page']}: {t['text'][:360]}...")
    return "\n".join(lines)


FILING_CATEGORIES = {
    "invoice": "Bills",
    "receipt": "Receipts",
    "contract": "Contracts",
    "bank": "Bank",
    "id": "IDs and Personal",
    "medical": "Medical",
    "certificate": "Certificates",
    "legal": "Legal",
    "news": "Articles",
    "general": "Other",
}

# Subcategory rules for bills/receipts, UAE-first, Arabic + English keywords.
FILING_SUBTYPES = [
    ("Utility", ["dewa", "ديوا", "addc", "sewa", "fewa", "كهرباء", "electricity", "الكهرباء", "مياه", "ماء", "water", "غاز", "gas", "هيئه كهرباء"]),
    ("Telecom", ["etisalat", "اتصالات", "du ", "virgin mobile", "e&", "stc", "mobily", "زين", "zain", "انترنت", "internet", "فايبر", "fiber", "جوال", "mobile plan"]),
    ("Rent", ["ايجار", "rent", "lease", "عقار", "إيجار", "ejari", "ايجاري", "tenancy"]),
    ("Government", ["حكوم", "رسوم", "مرور", "جوازات", "بلدي", "تأشيره", "visa fee", "government", "municipality", "salik", "سالك", "amer", "tas-heel", "تسهيل"]),
    ("Shopping", ["امازون", "amazon", "نون", "noon", "كارفور", "carrefour", "شرف", "sharaf", "متجر", "store", "shop"]),
]


def parse_date_ym(d: dict) -> str:
    """Best-effort YYYY-MM from the guessed date, else file mtime."""
    raw = normalize_digits(d.get("date_guess") or "")
    m = re.search(r"(20\d{2})[-/.](\d{1,2})", raw)
    if m:
        return f"{m.group(1)}-{int(m.group(2)):02d}"
    m = re.search(r"(\d{1,2})[-/.](\d{1,2})[-/.](20\d{2})", raw)
    if m:
        return f"{m.group(3)}-{int(m.group(2)):02d}"
    m = re.search(r"(jan|feb|mar|apr|may|jun|jul|aug|sep|oct|nov|dec)[a-z]*\s+(20\d{2})", raw, re.I)
    if m:
        months = {"jan": 1, "feb": 2, "mar": 3, "apr": 4, "may": 5, "jun": 6, "jul": 7, "aug": 8, "sep": 9, "oct": 10, "nov": 11, "dec": 12}
        return f"{m.group(2)}-{months[m.group(1).lower()[:3]]:02d}"
    try:
        return datetime.fromtimestamp(d.get("mtime") or 0).strftime("%Y-%m")
    except Exception:
        return datetime.now().strftime("%Y-%m")


def filing_subcategory(d: dict) -> str:
    doc_type = d.get("doc_type") or "general"
    hay = normalize_arabic(" ".join([
        str(d.get("company") or ""), str(d.get("title") or ""),
        str(d.get("summary") or ""), str(d.get("normalized_text") or "")[:2000],
    ]))
    if doc_type in ("invoice", "receipt"):
        for name, kws in FILING_SUBTYPES:
            # Word-boundary matching, not substring: "du" as a substring
            # matches inside "due_or_expiry_date" and mislabels everything.
            if any(contains_keyword(hay, kw) for kw in kws):
                return name
        return "General"
    # Other document types file by year.
    return parse_date_ym(d)[:4]


def filing_suggestion(d: dict, base: Optional[Path] = None, company_counts: Optional[dict] = None) -> dict:
    """Hierarchical filing plan: Category/Subcategory[/Company]/typecode_company_YYYY-MM.ext
    A company subfolder is created only when that company has 3+ documents
    (company_counts) - otherwise files sit flat in the subcategory and the
    filename carries the company. Naming pattern overridable with
    DOCWISE_FILING_PATTERN using {type} {company} {date} {title} tokens."""
    base = base or ARCHIVE
    doc_type = d.get("doc_type") or "general"
    category = FILING_CATEGORIES.get(doc_type, "Other")
    subcategory = filing_subcategory(d)
    company_raw = (d.get("company") or "").strip()
    company = normalize_filename(company_raw)[:40]
    ym = parse_date_ym(d)
    type_code = {"invoice": "bill", "receipt": "receipt", "contract": "contract", "bank": "bank",
                 "id": "id", "medical": "medical", "certificate": "certificate", "legal": "legal",
                 "news": "article", "general": "doc"}.get(doc_type, "doc")
    pattern = os.environ.get("DOCWISE_FILING_PATTERN", "{type}_{company}_{date}")
    stem = pattern.format(
        type=type_code,
        company=company or "unknown",
        date=ym,
        title=normalize_filename(str(d.get("title") or ""))[:40],
    )
    ext = d.get("file_ext") or Path(d.get("path", "file.pdf")).suffix
    filename = normalize_filename(stem)[:120] + ext
    # A provider/company subfolder only when it earns its place: business
    # document types AND at least 3 documents from that company. One-off
    # companies file flat - their name is already in the filename.
    company_level = ""
    if company and doc_type in ("invoice", "receipt", "bank", "contract"):
        if company_counts and company_counts.get(company, 0) >= 3:
            company_level = company[:30]
    parts = [category, subcategory] + ([company_level] if company_level else [])
    folder = base.joinpath(*parts)
    return {
        "category": category,
        "subcategory": subcategory,
        "company": company_raw or None,
        "folder": str(folder.resolve()),
        "filename": filename,
        "target_path": str((folder / filename).resolve()),
        "relative": "/".join(parts + [filename]),
    }


def normalize_filename(value: str) -> str:
    value = re.sub(r"[^A-Za-z0-9\u0600-\u06FF_.-]+", "_", value or "")
    value = re.sub(r"_+", "_", value).strip("_.-")
    return value or "document"


@app.post("/api/organize")
def organize(req: OrganizeRequest):
    with DB_LOCK, db() as conn:
        row = conn.execute("SELECT * FROM documents WHERE id=?", (req.document_id,)).fetchone()
    if not row:
        raise HTTPException(status_code=404, detail="Document not found")
    d = row_to_doc(row, include_text=True)
    source = Path(d["path"])
    if not source.exists():
        raise HTTPException(status_code=404, detail="Original file missing")
    suggestion = filing_suggestion(d)
    target = Path(suggestion["target_path"])
    target.parent.mkdir(parents=True, exist_ok=True)
    if target.exists():
        target = target.with_name(f"{target.stem}_{int(time.time())}{target.suffix}")
    if req.mode == "move":
        shutil.move(str(source), str(target))
        new_path = target
    else:
        shutil.copy2(source, target)
        new_path = target
    indexed = index_file(new_path, "archive", force=True)
    return {"ok": True, "mode": req.mode, "target": str(target), "indexed": indexed}


class FilingPlanRequest(BaseModel):
    ids: list[int] = []
    base: str = ""


class FilingApplyRequest(BaseModel):
    ids: list[int] = []
    base: str = ""
    mode: str = "copy"


def _filing_rows(ids: list[int]):
    with DB_LOCK, db() as conn:
        if ids:
            placeholders = ",".join("?" for _ in ids)
            return conn.execute(f"SELECT * FROM documents WHERE id IN ({placeholders})", ids).fetchall()
        return conn.execute("SELECT * FROM documents ORDER BY doc_type, company, date_guess").fetchall()


def _company_counts(rows) -> dict:
    counts: dict = {}
    for r in rows:
        if (r["doc_type"] or "") in ("invoice", "receipt", "bank", "contract"):
            c = normalize_filename((r["company"] or "").strip())[:40]
            if c:
                counts[c] = counts.get(c, 0) + 1
    return counts


def prune_empty_dirs(base: Path) -> int:
    """Remove empty directories left behind after moves (deepest first)."""
    removed = 0
    try:
        for p in sorted((d for d in base.rglob("*") if d.is_dir()), key=lambda x: len(str(x)), reverse=True):
            try:
                p.rmdir()
                removed += 1
            except OSError:
                pass
    except Exception:
        pass
    return removed


@app.post("/api/filing/plan")
def filing_plan(req: FilingPlanRequest):
    """Dry run: how every document WOULD be filed. Touches nothing."""
    base = Path(req.base).expanduser() if req.base.strip() else ARCHIVE
    items = []
    rows = _filing_rows(req.ids)
    counts = _company_counts(rows)
    for r in rows:
        d = row_to_doc(r, include_text=True)
        s = filing_suggestion(d, base=base, company_counts=counts)
        items.append({
            "id": d["id"], "title": d.get("title"), "current_path": d["path"],
            "file_exists": Path(d["path"]).exists(),
            "category": s["category"], "subcategory": s["subcategory"], "company": s["company"],
            "relative": s["relative"], "target_path": s["target_path"],
        })
    return {"base": str(base.resolve()), "count": len(items), "items": items}


def file_one_document(d: dict, base: Path, mode: str, company_counts: Optional[dict] = None) -> dict:
    """Copy/move one document into the filing tree. Never overwrites: name
    collisions get a numeric suffix, identical content in place is skipped.
    mode 'auto' = copy, plus a guard against re-filing files already inside
    the tree (prevents watcher loops)."""
    source = Path(d["path"])
    item = {"id": d["id"], "from": str(source)}
    try:
        if not source.exists():
            raise FileNotFoundError("original file missing")
        s = filing_suggestion(d, base=base, company_counts=company_counts)
        target = Path(s["target_path"])
        if source.resolve() == target.resolve():
            item.update({"status": "already-filed", "to": str(target)})
            return item
        if mode == "auto" and str(source.resolve()).lower().startswith(str(base.resolve()).lower()):
            item.update({"status": "inside-archive"})
            return item
        target.parent.mkdir(parents=True, exist_ok=True)
        sha = d.get("sha256") or file_hash(source)
        n = 2
        while target.exists() and file_hash(target) != sha:
            target = target.with_name(f"{Path(s['filename']).stem}-{n}{target.suffix}")
            n += 1
        if target.exists():
            item.update({"status": "already-exists", "to": str(target)})
        elif mode == "move":
            shutil.move(str(source), str(target))
            with DB_LOCK, db() as conn:
                conn.execute(
                    "UPDATE documents SET path=?, source_type='archive', updated_at=? WHERE id=?",
                    (str(target), now_iso(), d["id"]),
                )
            item.update({"status": "moved", "to": str(target)})
        else:
            shutil.copy2(str(source), str(target))
            item.update({"status": "copied", "to": str(target)})
    except Exception as exc:
        item.update({"status": "error", "error": str(exc)})
    return item


def auto_file_document(doc_id: int) -> Optional[dict]:
    base_env = os.environ.get("DOCWISE_AUTO_FILE_BASE", "")
    base = Path(base_env).expanduser() if base_env.strip() else ARCHIVE
    with DB_LOCK, db() as conn:
        row = conn.execute("SELECT * FROM documents WHERE id=?", (doc_id,)).fetchone()
        counts = None
        if row and (row["company"] or "").strip():
            cnt = conn.execute(
                "SELECT COUNT(*) c FROM documents WHERE company=? AND doc_type IN ('invoice','receipt','bank','contract')",
                (row["company"],),
            ).fetchone()["c"]
            counts = {normalize_filename(row["company"].strip())[:40]: cnt}
    if not row:
        return None
    d = row_to_doc(row, include_text=True)
    item = file_one_document(d, base, "auto", company_counts=counts)
    return item if item.get("status") in ("copied", "moved") else None


@app.post("/api/filing/apply")
def filing_apply(req: FilingApplyRequest):
    """Execute the filing plan: copy (default) or move files into the tree."""
    if req.mode not in ("copy", "move"):
        raise HTTPException(status_code=400, detail="mode must be copy or move")
    base = Path(req.base).expanduser() if req.base.strip() else ARCHIVE
    done = errors = 0
    items = []
    rows = _filing_rows(req.ids)
    counts = _company_counts(rows)
    for r in rows:
        d = row_to_doc(r, include_text=True)
        item = file_one_document(d, base, req.mode, company_counts=counts)
        if item.get("status") in ("copied", "moved"):
            done += 1
        elif item.get("status") == "error":
            errors += 1
        items.append(item)
    pruned = prune_empty_dirs(base) if req.mode == "move" else 0
    return {"ok": True, "mode": req.mode, "base": str(base.resolve()), "done": done, "errors": errors, "pruned_empty_folders": pruned, "items": items}


@app.post("/api/open/{doc_id}")
def open_file(doc_id: int):
    with DB_LOCK, db() as conn:
        row = conn.execute("SELECT path FROM documents WHERE id=?", (doc_id,)).fetchone()
    if not row:
        raise HTTPException(status_code=404, detail="Document not found")
    path = Path(row["path"])
    if not path.exists():
        raise HTTPException(status_code=404, detail="File not found")
    try:
        if hasattr(os, "startfile"):
            os.startfile(str(path))  # type: ignore[attr-defined]
        elif sys.platform == "darwin":
            subprocess.run(["open", str(path)], check=False)
        else:
            subprocess.run(["xdg-open", str(path)], check=False)
        return {"ok": True}
    except Exception as exc:
        return {"ok": False, "error": str(exc), "path": str(path)}


@app.get("/api/file/{doc_id}")
def get_file(doc_id: int):
    with DB_LOCK, db() as conn:
        row = conn.execute("SELECT path FROM documents WHERE id=?", (doc_id,)).fetchone()
    if not row:
        raise HTTPException(status_code=404, detail="Document not found")
    path = Path(row["path"])
    if not path.exists():
        raise HTTPException(status_code=404, detail="File not found")
    return FileResponse(path)


app.mount("/", StaticFiles(directory=STATIC, html=True), name="static")


if __name__ == "__main__":
    import uvicorn
    uvicorn.run("app:app", host="127.0.0.1", port=8120, reload=False)
